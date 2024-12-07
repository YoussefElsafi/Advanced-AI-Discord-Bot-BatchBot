# Files
from system.config import TOKEN, NAME, API_KEY, sys_security, gen_config, gen_config2, HUGGING_FACE_API, Image_Model, DEFAULT_MUSIC_MODEL, history_limit, limit_history, \
    show_time, history_channel_toggle, embed_colors, Object_Detection_Model, show_tokens_at_startup, fix_repeating_prompts, safe_search, ffmpeg_executable_path, tts_toggle, \
        vc_voice, VOICES, sync_voice_with_text, HISTORY_FILE, smart_recognition, show_invite_link_on_startup, safegen, discord_heartbeat_timeout, mod_channel_name, \
            preview_code_output, additional_details, model_name, preview_model_name, model_temperature, create_mod_channel, show_tokens, add_watermark_to_generated_image, \
                show_safety_settings_on_startup, Dangerous, Harassment, Hate_Speech, Sexually_Explicit, Dangerous_Content, vc_AI, web_search, GOOGLE_CUSTOM_SEARCH_API_KEY, \
                    GOOGLE_PROJECT_SEARCH_ENGINE_ID, advanced_model
from system.instructions.instruction import ins, video_ins, file_ins, insV, insV2, fix_mem_ins, cool_ins
from system.instructions.instruction_ru import ru_ins, ru_video_ins, ru_file_ins, ru_insV, ru_insV2, ru_fix_mem_ins
from system.instructions.instruction_eg import eg_ar_ins, eg_fix_mem_ins
from system.instructions.instruction_fr import fr_ins, fr_video_ins, fr_file_ins, fr_insV, fr_insV2, fr_fix_mem_ins
from system.instructions.instruction_es import es_ins, es_video_ins, es_file_ins, es_fix_mem_ins
from system.instructions.instruction_de import de_ins, de_video_ins, de_file_ins, de_insV, de_insV2, de_fix_mem_ins
from system.instructions.instruction_ar import ins_ar, ar_fix_mem_ins
from system.instructions.instruction_tutor_mode import tutor_ins
import system.check_tokens as check
from system.check_tokens import tokens

# Libraries
import discord
from discord.ext import commands
import google.generativeai as genai
import json
import os
import requests
from PIL import Image
from colorama import Fore, Style
import asyncio
import logging
import random
import time
import httpx
from discord.utils import get
import io
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api import TranscriptsDisabled
import urllib.parse as urlparse
import re
from urllib.parse import urlparse, parse_qs
import inspect
import docx
import pptx
import openpyxl
import datetime
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import edge_tts
import re
import shutil
from bs4 import BeautifulSoup

discord_verified, gemini_api_key_verified, hugging_verified, google_search_api_key_verified, google_search_project_id_verified = tokens()

if not discord_verified:
    exit()

if google_search_api_key_verified and google_search_project_id_verified:
    google_search_api_verified = True
else:
    google_search_api_verified = False

# Set up the bot with the correct prefix and intents
intents = discord.Intents.default()
intents.members = True
intents.presences = True
intents.message_content = True
ffmpeg_path = ffmpeg_executable_path

bot = commands.Bot(command_prefix="/", intents=intents, heartbeat_timeout=discord_heartbeat_timeout)

dev_DEBUG = False
Model_Debug = False

# Ensure the log directory exists before configuring logging
log_dir = "system/log"
os.makedirs(log_dir, exist_ok=True)
if os.path.exists(log_dir):
    timestamp = time.strftime('%Y-%m-%d_%H-%M-%S')  # Replace colons with underscores
    logging.basicConfig(filename=f"system/log/{timestamp}.log", level=logging.DEBUG,
                        format='%(asctime)s - %(levelname)s - %(filename)s - %(lineno)d - %(message)s')

SEARCH_API_KEY = GOOGLE_CUSTOM_SEARCH_API_KEY
CX = GOOGLE_PROJECT_SEARCH_ENGINE_ID

def fetch_code_and_content(url):
    """
    Fetches detailed content and code snippets from a given URL.
    Tries to scrape the most relevant content, including code blocks.
    """
    try:
        response = requests.get(url, timeout=10)  # Added timeout to avoid long delays
        if response.status_code == 200:
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Try to find main content
            main_content = (
                soup.find('article') or soup.find('section') or
                soup.find('main') or soup.find('div', {'id': 'content'}) or
                soup.find('div', {'class': 'main-content'})
            )

            # Try to find code blocks
            code_blocks = soup.find_all(['pre', 'code'])
            code_snippets = [
                code.get_text(strip=True)[:200] for code in code_blocks
            ]  # Limiting to first 200 characters per block
            
            # Extract paragraphs if no main content is found
            content = ""
            if main_content:
                paragraphs = main_content.find_all('p')
                content = ' '.join(p.get_text() for p in paragraphs[:5])  # First 5 paragraphs
            
            if code_snippets:
                content += '\nCode/Content Snippets:\n' + '\n'.join(code_snippets)
            
            return content.strip() if content.strip() else "Unable to retrieve meaningful content."
        return f"Error fetching page: {response.status_code}"
    except Exception as e:
        return f"Error fetching content from URL: {str(e)}"

def search_google(query, site=None, num_results=5, safe_search=safe_search):
    """
    Performs a Google Custom Search and fetches detailed content (including code) from the top search results.

    Args:
        query (str): The search query.
        site (str, optional): The site to restrict the search to (e.g., "https://www.youtube.com/"). Default is None.
        num_results (int, optional): The number of results to fetch. Default is 5.
        safe_search (bool, optional): Whether to enable SafeSearch. Default is True.

    Returns:
        list: A list of search result dictionaries or an empty list if no results are found.
    """
    safe_search_set = "active" if safe_search else "off"

    # Set up API request
    url = "https://www.googleapis.com/customsearch/v1"
    params = {
        "key": SEARCH_API_KEY,
        "cx": CX,
        "q": query,
        "num": min(num_results, 10),  # Google Custom Search API allows up to 10 results per request
        "safe": safe_search_set,
    }
    if site:
        params["siteSearch"] = site

    # Perform the request
    try:
        response = requests.get(url, params=params, timeout=10)
        response.raise_for_status()
        results = response.json()

        if "items" in results:
            return [
                {
                    "title": item.get("title"),
                    "link": item.get("link"),
                    "snippet": item.get("snippet"),
                }
                for item in results["items"]
            ]
        else:
            print("No results found.")
            return "No results found."
    except requests.exceptions.RequestException as e:
        print(f"Error during search request: {e}")
        return "Error Searching the web."


if not os.path.exists('system/data'):
    os.makedirs('system/data') 
if not os.path.exists('system/RAM'):
    os.makedirs('system/RAM') 

async def send_message(channel, message, max_length=1999):
    """
    Split a message into multiple chunks and send them to the given channel.
    
    The message is split into chunks of up to max_length characters. The message
    is split at newline characters and the chunks are then sent to the channel
    one by one. If the message is too long, it is split into multiple chunks and
    sent separately.
    """
    lines = message.splitlines()
    chunks = []
    current_chunk = ""
    
    for line in lines:
        if len(current_chunk) + len(line) + 1 > max_length:
            chunks.append(current_chunk.strip())
            current_chunk = line
        else:
            if current_chunk:
                current_chunk += "\n"
            current_chunk += line
    if current_chunk:
        chunks.append(current_chunk.strip())
    
    for part in chunks:
        await channel.send(part)

# Function to load conversation history from file
def load_history():
    if os.path.exists(HISTORY_FILE):
        with open(HISTORY_FILE, 'r') as file:
            return json.load(file)
    return {}

# Initialize conversation history
conversation_history = load_history()

# Function to save conversation history to file
def save_history():
    # Create the directory if it does not exist
    os.makedirs(os.path.dirname(HISTORY_FILE), exist_ok=True)
    with open(HISTORY_FILE, 'w') as file:
        json.dump(conversation_history, file, indent=4)

# Function to add a message to the conversation history
def add_to_history(member_name, message, channel_name=None):
    """Adds a message to the conversation history."""
    timestamp = time.strftime('%Y-%m-%d %H:%M:%S')

    # Get context object ('message', 'ctx.message', etc.) from calling function
    frame = inspect.currentframe().f_back
    args, _, _, values = inspect.getargvalues(frame)
    context_obj = values.get('message', None)
    if context_obj is None:
        context_obj = values.get('ctx', None)
        if context_obj is not None:
            context_obj = context_obj.message

    # Use provided channel_name or get it from context
    if channel_name is None:  
        if history_channel_toggle and context_obj is not None:
            user_id = context_obj.channel.name
        else:
            user_id = "Conversation"
    else:
        user_id = channel_name  # Use the provided channel_name

    if user_id not in conversation_history:
        conversation_history[user_id] = []

    if show_time:
        conversation_history[user_id].append(f"{timestamp} - {member_name}: {message}")
    else:
        conversation_history[user_id].append(f"{member_name}: {message}")

    # Truncate history if limit_history is True 
    if limit_history and len(conversation_history[user_id]) > history_limit:
        conversation_history[user_id] = conversation_history[user_id][-history_limit:]

    save_history() 

def add_to_history_bot(member_name, message, channel_name=None):
    """Adds a bot message to the conversation history."""
    timestamp = time.strftime('%Y-%m-%d %H:%M:%S')

    # Get context object ('message', 'ctx.message', etc.) from calling function
    frame = inspect.currentframe().f_back
    args, _, _, values = inspect.getargvalues(frame)
    context_obj = values.get('message', None)
    if context_obj is None:
        context_obj = values.get('ctx', None)
        if context_obj is not None:
            context_obj = context_obj.message

    # Use provided channel_name or get it from context
    if channel_name is None:  
        if history_channel_toggle and context_obj is not None:
            user_id = context_obj.channel.name
        else:
            user_id = "Conversation"
    else:
        user_id = channel_name  # Use the provided channel_name

    if user_id not in conversation_history:
        conversation_history[user_id] = []

    if show_time:
        conversation_history[user_id].append(f"{timestamp} - {member_name}{message}")
    else:
        conversation_history[user_id].append(f"{member_name}{message}")
    save_history()
if not history_channel_toggle:
    add_to_history("System", "You have been rebooted!")

async def unnecessary_error(e):
    error_message = str(e)
    if "500" in error_message:
        return True
    else:
        return False

async def debug_error(e, message, channel):
    error_message = str(e)
    if "500" in error_message:
        logging.error(f"Google Internal Error while {message}: {e}")
        add_to_history("Google Internal Error", "💥 An Internal Error has occured, Retrying...")
        print(f"Google Error while {message}: {e}")
    elif "503" in error_message:
        logging.warning(f"Error (Temporarily overloaded or down service) While {message}: {e}")
        add_to_history("Error", "The service may be temporarily overloaded or down. Please try again later.")
        print(f"Error (Temporarily overloaded or down service) While {message}: {e}")
        await channel.send("⏳ The service may be temporarily overloaded or down. Please try again later.")
    elif "403" in error_message:
        logging.error(f"Error (API Key Denied Permissions) While {message}: {e}")
        add_to_history("Error", "Your API key doesn't have the required permissions.")
        await channel.send("🔒 Your API key has denied permissions.")
        print(f"Error (API Key Denied Permissions) While {message}: {e}")
    elif "504" in error_message:
        logging.warning(f"Error (Service Unable to finish processing within the deadline) While {message}: {e}")
        add_to_history("Error", "The service is unable to finish processing within the deadline.")
        await channel.send("⏳ The service is unable to finish processing within the deadline.")
        print(f"Error (Service Unable to finnish processing within the deadline) While {message}: {e}")
    elif "429" in error_message:
        logging.warning(f"Error (Service rate limited) While {message}: {e}")
        add_to_history("Error", "The service is being rate limited.")
        await channel.send("🚫 You've exceeded the rate limit, Please try again later.")
        print(f"Error (Service rate limited) While {message}: {e}")
    else:
        logging.error(f"An Error occured while {message}: {e}")
        print(f"An Error occured while {message}: {e}")
        add_to_history("Error", f"Error occurred while {message}: {error_message}")
        await channel.send("🚫 Uh oh! Something went wrong. We couldn't complete the request. Please try again.")

# Utility functions
def save_search(query, result):
    with open('system/data/saved-searches.py', 'a') as f:
        f.write(f'{query}: {result} |\n')

def save_memory(query, result):
    """Saves memory to a JSON file."""
    try:
        # Load existing memory if it exists
        with open('system/data/core-memory.json', 'r') as f:
            memory = json.load(f)
    except FileNotFoundError:
        memory = {}

    # Add the new memory entry
    memory[query] = result

    # Save the updated memory
    with open('system/data/core-memory.json', 'w') as f:
        json.dump(memory, f, indent=4)

def load_memory(query=None):
    """Loads memory from a JSON file."""
    try:
        with open('system/data/core-memory.json', 'r') as f:
            memory = json.load(f)
            if query:
                return memory.get(query)
            else:
                return memory
    except FileNotFoundError:
        return {}

def get_conversation_history(ctx=None): 
    """Gets the conversation history based on history_channel_toggle."""
    if history_channel_toggle and ctx is not None:
        user_id = ctx.channel.name
    else:
        user_id = "Conversation"
    return "\n".join(conversation_history.get(user_id, []))

api_key = f"{API_KEY}"
name = f"{NAME}"

check.tokens()

if show_tokens_at_startup:
    print(" ")
    print(f"{Fore.WHITE + Style.BRIGHT + Style.DIM}API KEY:{Style.RESET_ALL} {Fore.MAGENTA + Style.BRIGHT}{api_key}{Style.RESET_ALL}")
    print(Fore.RED + Style.BRIGHT + "__________________________________________________________________________________")
    print(" ")
    print(f"{Fore.WHITE + Style.BRIGHT + Style.DIM}BOT TOKEN:{Style.RESET_ALL} {Fore.BLUE + Style.BRIGHT}{TOKEN}{Style.RESET_ALL}")
    print(Fore.RED + Style.BRIGHT + "__________________________________________________________________________________")
    print(" ")
    print(f"{Fore.WHITE + Style.BRIGHT + Style.DIM}HUGGING FACE API KEY:{Style.RESET_ALL} {Fore.YELLOW + Style.BRIGHT}{HUGGING_FACE_API}{Style.RESET_ALL}")
    print(" ")

# Global variable to store the member's custom name
member_custom_name = {}

@bot.tree.command(name="name", description="Change your custom name")
async def change_name(interaction: discord.Interaction, new_name: str):
    global member_custom_name
    if not new_name:  # Check for empty string
        await interaction.response.send_message("Please provide a name.", ephemeral=True)
    else:
        member_custom_name[interaction.user.id] = new_name
        await interaction.response.send_message(f"Your name has been changed to {new_name}.", ephemeral=True)

print(f"Defualt Model: {preview_model_name}")
if show_safety_settings_on_startup:
    print(f"{Fore.GREEN + Style.BRIGHT}Safety Settings:{Style.RESET_ALL}")
    print(f"{Fore.CYAN + Style.BRIGHT}Dangerous Safety Set: {Fore.RESET}{Dangerous}")
    print(f"{Fore.CYAN + Style.BRIGHT}Harassment Safety Set: {Fore.RESET}{Harassment}")
    print(f"{Fore.CYAN + Style.BRIGHT}Hate Speech Safety Set: {Fore.RESET}{Hate_Speech}")
    print(f"{Fore.CYAN + Style.BRIGHT}Sexually Explicit Safety Set: {Fore.RESET}{Sexually_Explicit}")
    print(f"{Fore.CYAN + Style.BRIGHT}Dangerous Content Safety Set: {Fore.RESET}{Dangerous_Content}")
    print("")
genai_model = model_name

if fix_repeating_prompts:
    ins = f"{ins}\n{fix_mem_ins}"

default_lang = "en"

ins = f"{ins}\n{cool_ins}\nAnd if someone asks you what is the best model, tell them its Google Gemini Experimental 1206 because it outperforms all other models. even Google Gemini 1.5 Pro and Google Gemini Ultra and all openai's models, and it can handle alot and every complex task! It's an EXTREMELY advanced model! and ranks #1 on coding, #2 math, #1 on creative writing and #1 on Hard Prompts, Also if a user sends you a file attachment, dont use code to find the file\nalso! you can search youtube videos using `/search*yt [query]` and the same with searching the web! `/search [query]`! its that easy! so if a user asks to search something on youtube/the web, then do this! and instead of making codes to search the web, just use `/search [query]` to search the web! you can search the web with /search [query]! and same with youtube! /search*yt [query]!, and when you try doing text links like [a link](https://google.com), you cant do '[https://www.youtube.com/watch?v=ENrzD9HAZK4](https://www.youtube.com/watch?v=ENrzD9HAZK4)' but you can do '[a](https://www.youtube.com/watch?v=ENrzD9HAZK4)', so you cannot link a link but you can link a text. so if you want to give the link and visible without hiding it with a text, just put the link and without any () and []!\n"                        
if genai_model == "learnlm-1.5-pro-experimental":
    if fix_repeating_prompts:
        tutor_ins = (f"{tutor_ins}\n{fix_mem_ins}\n{cool_ins}\nAnd if someone asks you what is the best model, tell them its Google Gemini Experimental 1206 because it outperforms all other models. even Google Gemini 1.5 Pro and Google Gemini Ultra and all openai's models, and it can handle alot and every complex task! It's an EXTREMELY advanced model! and ranks #1 on coding, #2 math, #1 on creative writing and #1 on Hard Prompts, Also if a user sends you a file attachment, dont use code to find the file\nalso! you can search youtube videos using `/search*yt [query]` and the same with searching the web! `/search [query]`! its that easy! so if a user asks to search something on youtube/the web, then do this! and instead of making codes to search the web, just use `/search [query]` to search the web! you can search the web with /search [query]! and same with youtube! /search*yt [query]!, and when you try doing text links like [a link](https://google.com), you cant do '[https://www.youtube.com/watch?v=ENrzD9HAZK4](https://www.youtube.com/watch?v=ENrzD9HAZK4)' but you can do '[a](https://www.youtube.com/watch?v=ENrzD9HAZK4)', so you cannot link a link but you can link a text. so if you want to give the link and visible without hiding it with a text, just put the link and without any () and []!",),
    else:
        tutor_ins = (f"{tutor_ins}\n{cool_ins}\nAnd if someone asks you what is the best model, tell them its Google Gemini Experimental 1206 because it outperforms all other models. even Google Gemini 1.5 Pro and Google Gemini Ultra and all openai's models, and it can handle alot and every complex task! It's an EXTREMELY advanced model! and ranks #1 on coding, #2 math, #1 on creative writing and #1 on Hard Prompts, Also if a user sends you a file attachment, dont use code to find the file\nalso! you can search youtube videos using `/search*yt [query]` and the same with searching the web! `/search [query]`! its that easy! so if a user asks to search something on youtube/the web, then do this! and instead of making codes to search the web, just use `/search [query]` to search the web! you can search the web with /search [query]! and same with youtube! /search*yt [query]!, and when you try doing text links like [a link](https://google.com), you cant do '[https://www.youtube.com/watch?v=ENrzD9HAZK4](https://www.youtube.com/watch?v=ENrzD9HAZK4)' but you can do '[a](https://www.youtube.com/watch?v=ENrzD9HAZK4)', so you cannot link a link but you can link a text. so if you want to give the link and visible without hiding it with a text, just put the link and without any () and []!",),

# Configure the Google Generative AI
genai.configure(api_key=f"{api_key}")
# The core model
model = genai.GenerativeModel( 
    model_name=genai_model,
    generation_config=gen_config,
    system_instruction = ins if genai_model != "learnlm-1.5-pro-experimental" else tutor_ins,
    safety_settings=sys_security,
    tools='code_execution' if preview_code_output else None
)

# Other Models...
model_flash = genai.GenerativeModel( 
    model_name="gemini-1.5-flash",
    generation_config=gen_config,
    system_instruction=(ins),
    safety_settings=sys_security
)
model_pro = genai.GenerativeModel( 
    model_name="gemini-1.5-pro-latest",
    generation_config=gen_config,
    system_instruction=(insV),
    safety_settings=sys_security
)
model_V = genai.GenerativeModel( 
    model_name=advanced_model,
    generation_config=gen_config,
    system_instruction=(insV),
    safety_settings=sys_security
)
model_V2 = genai.GenerativeModel( 
    model_name="gemini-1.5-flash",
    generation_config=gen_config,
    system_instruction=(insV),
    safety_settings=sys_security
)
model_V3 = genai.GenerativeModel( 
    model_name="gemini-1.5-flash",
    generation_config=gen_config,
    system_instruction=(insV2),
    safety_settings=sys_security
)
model3 = genai.GenerativeModel( 
    model_name="gemini-1.5-flash",
    generation_config=gen_config2,
    system_instruction=("MAX LENGTH IS 80 WORDS"),
    safety_settings=sys_security
)
model_name = genai.GenerativeModel(
  model_name="gemini-1.5-flash",
  generation_config=gen_config,
  system_instruction="you are only an memory name generator engine, generate memory names only as the memory prompted and dont say anything else, the system will tell you what to generate, only generate 1 name and dont make it too long and make it silly, and DONT say `/n:` and i used / instead of the other one because it is gonna break the system",
  safety_settings=sys_security
)
model_vid = genai.GenerativeModel(
    model_name="gemini-1.5-pro-latest",
    generation_config=gen_config,
    system_instruction=(video_ins),
    safety_settings=sys_security
)
model_vid_a = genai.GenerativeModel(
    model_name=advanced_model,
    generation_config=gen_config,
    system_instruction=(video_ins),
    safety_settings=sys_security
)
model_vid_flash = genai.GenerativeModel(
    model_name="gemini-1.5-flash",
    generation_config=gen_config,
    system_instruction=(video_ins),
    safety_settings=sys_security
)
model_file = genai.GenerativeModel(
    model_name="gemini-1.5-pro-latest",
    generation_config=gen_config,
    system_instruction=(file_ins),
    safety_settings=sys_security
)
model_file_a = genai.GenerativeModel(
    model_name=advanced_model,
    generation_config=gen_config,
    system_instruction=(file_ins),
    safety_settings=sys_security
)
model_file_flash = genai.GenerativeModel(
    model_name="gemini-1.5-flash",
    generation_config=gen_config,
    system_instruction=(file_ins),
    safety_settings=sys_security
)
model_object = genai.GenerativeModel(
    model_name="gemini-1.5-flash-latest",
    generation_config=gen_config,
    system_instruction="Your only propose is to get the details that the user sent to you and you convert them into human talk only and nothing else, example: 'User: [{'score': 0.9994643330574036, 'label': 'sports ball', 'box': {'xmin': 95, 'ymin': 444, 'xmax': 172, 'ymax': 515}}, {'score': 0.810539960861206, 'label': 'person', 'box': {'xmin': 113, 'ymin': 15, 'xmax': 471, 'ymax': 414}}, {'score': 0.7840690612792969, 'label': 'person', 'box': {'xmin': 537, 'ymin': 35, 'xmax': 643, 'ymax': 241}}, {'score': 0.9249405860900879, 'label': 'person', 'box': {'xmin': 109, 'ymin': 14, 'xmax': 497, 'ymax': 528}}, {'score': 0.9990099668502808, 'label': 'person', 'box': {'xmin': 0, 'ymin': 47, 'xmax': 160, 'ymax': 373}}, {'score': 0.8631113767623901, 'label': 'person', 'box': {'xmin': 110, 'ymin': 13, 'xmax': 558, 'ymax': 528}}, {'score': 0.9433853626251221, 'label': 'person', 'box': {'xmin': 537, 'ymin': 34, 'xmax': 643, 'ymax': 310}}, {'score': 0.6196897625923157, 'label': 'person', 'box': {'xmin': 715, 'ymin': 160, 'xmax': 770, 'ymax': 231}}, {'score': 0.5696023106575012, 'label': 'person', 'box': {'xmin': 777, 'ymin': 170, 'xmax': 800, 'ymax': 221}}, {'score': 0.9989137649536133, 'label': 'person', 'box': {'xmin': 423, 'ymin': 67, 'xmax': 638, 'ymax': 493}}] | You: '- There's a sports ball near the bottom middle.\n- There are a few people in the image.\n- One person is on the left side.\n- A couple of people are in the center and middle-right.\n- There are a couple of possible people on the right, but the AI isn't as sure about them. \n' and you **MUST** use - at the start like in the example and only say the stuff that the user sent you and not anything else",
    safety_settings=sys_security
)

# Load existing conversation history from file
try:
    with open(HISTORY_FILE, 'r') as file:
        conversation_history = json.load(file)
except FileNotFoundError:
    conversation_history = {}

@bot.event
async def on_ready():
    await bot.add_cog(VoiceListener(bot))  # Await add_cog here
    print(f"Successfully Logged in as: {NAME}!")
    print("Bot is online! Type /help for a list of commands.")
    bot_invite_link = discord.utils.oauth_url(
        bot.user.id,
        permissions=discord.Permissions(),
        scopes=("bot", "applications.commands")
    )
    if show_invite_link_on_startup:
        print(f"Invite link: {bot_invite_link}")
    try:
        synced = await bot.tree.sync()
        if len(synced) > 1:
            print(f"Synced {len(synced)} commands")
        else:
            print(f"Synced {len(synced)} command")
    except Exception as e:
        print(f"{Fore.RED + Style.BRIGHT}Error:{Style.RESET_ALL} {e}")
        quit()
    print(Fore.WHITE + Style.BRIGHT + "__________________________________________________________________________________" + Style.RESET_ALL)
    print(" ")
    print(f"{Fore.MAGENTA + Style.BRIGHT}{NAME}'s Console:{Style.RESET_ALL}")
    print(" ")

EN_video_ins = video_ins
EN_insV = insV
EN_file_ins = file_ins
EN_insV2 = insV2
EN_ins = ins
        
# Start Gemini Chats //:
chat_session = model.start_chat(history=[])
chat_session_flash = model_flash.start_chat(history=[])

@bot.tree.command(name="report", description="Report a bug, issue or a user")
async def report(interaction: discord.Interaction, report: str):
    await interaction.response.defer()
    # Prepare the report entry
    user = interaction.user
    member_name = user.display_name
    report_entry = (
        "----------------------------------------------------------------------------------\n"
        f"Username: {user.name}#{user.discriminator} | Name: {member_name} (ID: {user.id})\n"
        f"Report: {report}\n"
        "----------------------------------------------------------------------------------\n\n"
    )

    # Path to the report file
    report_file_path = "system/data/reports.txt"

    # Write the report entry to the file
    with open(report_file_path, "a") as file:
        file.write(report_entry)

    add_to_history(member_name, f"System: {member_name} sent a report! `{report}`")
    await interaction.followup.send(f"Thank you for your report, {member_name}. `{report}` It has been logged.")

@bot.tree.command(name="feedback", description="Provide feedback or suggestions")
async def feedback(interaction: discord.Interaction, feedback: str):
    await interaction.response.defer()
    # Prepare the feedback entry
    user = interaction.user
    member_name = user.display_name
    feedback_entry = (
        "----------------------------------------------------------------------------------\n"
        f"Username: {user.name}#{user.discriminator} | Name: {member_name} (ID: {user.id})\n"
        f"Feedback: {feedback}\n"
        "----------------------------------------------------------------------------------\n\n"
    )

    # Path to the feedback file
    feedback_file_path = "system/data/feedback.txt"

    # Write the feedback entry to the file
    with open(feedback_file_path, "a") as file:
        file.write(feedback_entry)

    add_to_history(member_name, f"System: {member_name} sent feedback! `{feedback}`")
    await interaction.followup.send(f"Thank you for your feedback, {member_name}. `{feedback}` has been logged!")

# Function to check if the URL is a YouTube URL
def is_youtube_url(url):
    if url is None:
        return False
    youtube_regex = (
        r'(https?://)?(www\.)?'
        r'(youtube|youtu|youtube-nocookie)\.(com|be)/'
        r'(watch\?v=|embed/|v/|.+\?v=)?([^&=%\?]{11})'
    )
    return re.match(youtube_regex, url) is not None

# Function to extract video ID from a YouTube URL
def get_video_id(url):
    parsed_url = urlparse(url)
    if "youtube.com" in parsed_url.netloc:
        video_id = parse_qs(parsed_url.query).get('v')
        return video_id[0] if video_id else None
    elif "youtu.be" in parsed_url.netloc:
        return parsed_url.path[1:] if parsed_url.path else None
    return None

# Function to get the transcript from a YouTube video ID
def get_transcript_from_video_id(video_id):
    try:
        transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
        return ' '.join([i['text'] for i in transcript_list])
    except (KeyError, TranscriptsDisabled):
        return "Error retrieving transcript from YouTube video ID"

# Function to handle YouTube URLs, retrieve transcripts, and send them to the channel
async def handle_youtube_url(url, channel, prompt=None):
    """Handles YouTube URLs, retrieves transcripts, and sends them to the channel."""
    try:
        if not is_youtube_url(url):
            await channel.send("Invalid YouTube URL.")
            return

        video_id = get_video_id(url)
        if not video_id:
            await channel.send("Unable to extract video ID from URL.")
            return

        transcript = get_transcript_from_video_id(video_id)
        if "Error" in transcript:
            await channel.send(transcript)
            add_to_history("System", f"Error retrieving transcript: {transcript}")
        else:
            return transcript

    except Exception as e:
        await channel.send(f"An error occurred: {str(e)}")
        add_to_history("System", f"Error occurred: {str(e)}")

@bot.event
async def on_command_error(ctx, error):
    if isinstance(error, commands.CommandNotFound):
        await ctx.send(f"{error}")
        add_to_history("System", error)
    else:
        print(error)
        await ctx.send(f"An error occurred: {error}")
        raise error
    
# Constants
USER_SETTINGS_PATH = 'system/user-settings'

# Ensure the user settings directory exists
os.makedirs(USER_SETTINGS_PATH, exist_ok=True)

def get_user_settings(username):
    from system.config import model_name, preview_model_name
    user_file = os.path.join(USER_SETTINGS_PATH, f"{username}.json")
    default_settings = {
        'model': model_name,  # Use the model *name* string
        'model_name': preview_model_name # Display name
    }

    if not os.path.exists(user_file):
        print(f"Settings file not found for {username}. Creating...")
        try:
            with open(user_file, 'w') as file:
                json.dump(default_settings, file, indent=4)
            print(f"Settings file created for {username} with defaults.")
            with open(user_file, 'r') as file:
                settings = json.load(file)
            return settings
        except Exception as e:
            print(f"Error creating settings file: {e}")
            return default_settings  # Return defaults even if file creation fails

    else:  # Load from file if it exists
        try:
            with open(user_file, 'r') as file:
                settings = json.load(file)
            return settings

        except json.JSONDecodeError:
            print(f"Corrupted settings file for {username}. Recreating...")
            try:  # Try to recreate the file
                os.remove(user_file)
                with open(user_file, 'w') as file:
                    json.dump(default_settings, file, indent=4)
                return default_settings
            except Exception as e:  # Handle recreation errors
                print(f"Error recreating settings file: {e}")
                return default_settings # Defaults if can't recreate

        except Exception as e:
            print(f"Error loading settings file: {e}")
            return default_settings # Return defaults on unexpected error



def set_user_model(username, model):
    """Update the user's selected model."""
    user_file = os.path.join(USER_SETTINGS_PATH, f"{username}.json")
    user_settings = get_user_settings(username)
    user_settings['model'] = model

    model_name_mapping = {
        'gemini-1.5-pro': 'Gemini 1.5 Pro',
        'gemini-1.5-flash': 'Gemini 1.5 Flash',
        'gemini-1.5-flash-latest': 'Gemini 1.5 Flash Latest',
        'gemini-1.5-flash-8b': 'Gemini 1.5 Flash 8B',
        'gemini-1.5-pro-latest': 'Gemini 1.5 Pro Latest',
        'gemini-1.5-pro-002': 'Gemini 1.5 Pro 002',
        'gemini-1.5-flash-002': 'Gemini 1.5 Flash 002',
        'learnlm-1.5-pro-experimental': 'LearnLM 1.5 Pro',
        'gemini-exp-1114': 'Gemini Experimental 1114',
        'gemini-exp-1121': 'Gemini Experimental 1121',
        'gemini-exp-1206': 'Gemini Experimental 1206',
    }
    user_settings['model_name'] = model_name_mapping.get(model, model) # Default to the model string if not found

    with open(user_file, 'w') as file:
        json.dump(user_settings, file, indent=4)

def add_watermark(input_image_path, output_image_path):
    try:
        watermark_image_path = 'system/assets/watermark.png'
        """
        Adds an image watermark to the bottom-right corner of an image.
        
        Args:
            input_image_path (str): Path to the input image.
            output_image_path (str): Path to save the watermarked image.
            watermark_image_path (str): Path to the watermark image.
        """
        # Open the original image
        with Image.open(input_image_path) as img:
            # Open the watermark image

            with Image.open(watermark_image_path) as watermark:
                # Resize watermark if it is too large
                watermark_width, watermark_height = watermark.size
                max_width = int(img.size[0] * 0.25)  # Max width 25% of the input image
                if watermark_width > max_width:
                    watermark_ratio = max_width / watermark_width
                    watermark = watermark.resize(
                        (int(watermark_width * watermark_ratio), int(watermark_height * watermark_ratio)),
                        Image.Resampling.LANCZOS  # Use LANCZOS for high-quality downsampling
                    )

                # Calculate position for the watermark (bottom-left corner)
                x = 10  # 10px margin from the edge
                y = img.size[1] - watermark.height - 10
                # Paste the watermark onto the original image
                img.paste(watermark, (x, y), watermark.convert('RGBA').split()[3])  # Use alpha channel as mask

                # Save the watermarked image
                watermark_path = "system/assets"
                os.makedirs(os.path.dirname(watermark_path), exist_ok=True)
                img.save(output_image_path, "PNG")
    except FileNotFoundError:
        print("Error: The watermark image or input image not found.")
        
    except OSError:
        print("Error: Unable to save the image.")
    except Exception as e:
        print(f"Error adding a watermark: {e}")

async def handle_image_attachment(attachments, channel, prompt=None, message=None):
    """Handles the image attachment processing and deletion for multiple images."""
    
    image_files = []  # List to hold processed image paths

    for i, attachment in enumerate(attachments):
        file_extension = attachment.filename.split('.')[-1].lower()
        if file_extension == 'jpg':
            file_extension = 'jpeg'  # Rename 'jpg' to 'jpeg'
        
        # Generate a unique file name using a counter or timestamp
        unique_filename = f'image_{i}_{int(time.time())}.{file_extension}'
        file_path = os.path.join('system/RAM/read-img', unique_filename)
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        
        try:
            img_data = await attachment.read()
            with open(file_path, 'wb') as file:
                file.write(img_data)
            
            if file_extension == 'jpeg':
                img = Image.open(file_path).convert('RGB')
            else:
                img = Image.open(file_path)

            # Add image to the list of processed images
            image_files.append(img)
        
        except Exception as e:
            await channel.send(f"Error reading attachment {attachment.filename}: {str(e)}")
            print(f"Error reading attachment: {e}")
            return
    
    if image_files:
        
        
        # Prepare the images for the model and generate content
        while True:
            try:
                async with message.channel.typing():
                    print("DEBUG: Starting image processing for model input.")
                    
                    # Convert images to byte format for the model
                    image_file_objects = []
                    for img in image_files:
                        buffered = io.BytesIO()
                        img.save(buffered, format="PNG")
                        img_bytes = buffered.getvalue()
                        print(f"DEBUG: Processed image size in bytes: {len(img_bytes)}")
                        
                        image_file_objects.append({
                            'mime_type': 'image/png',  # or 'image/jpeg'
                            'data': img_bytes
                        })
                    print(f"DEBUG: Total number of images processed: {len(image_file_objects)}")
                    
                    # Retrieve the display name
                    display_name = member_custom_name.get(message.author.id, message.author.display_name)
                    member_name = display_name

                    # Fetch conversation history
                    model_conversation_history = get_conversation_history(message)
                    # Prepare the prompt with conversation history
                    mprompt = f"Conversation History Memory: |{model_conversation_history}|\n{member_name}: {prompt}"
                    response_text = ""  # Initialization
                    if not gemini_api_key_verified:
                        await message.channel.send("Please set up the API key at [Google AI Studio](https://aistudio.google.com/apikey)")
                        await message.add_reaction('🔑')
                        return
                    response = chat_session.send_message(image_file_objects + [mprompt])
                    response_text = response.text.strip()  # Assigning response_text here
                    add_to_history_bot("", f"{response_text}")

                    if response_text.startswith("/img"):
                        # Extract the text after "/img"
                        text_after_command = response_text[len("/img"):].strip() # //img

                        if text_after_command:
                            # Generate the text after "/img"
                            prompt_response_text = text_after_command
                            add_to_history_bot("", f"/img {prompt_response_text}")
                        else:
                            history = get_conversation_history(message) # Use the function to get history
                            full_prompt = f"{history}\nVisualizer: What image do you want to generate?: "
                            response = model.generate_content(full_prompt)
                            prompt_response_text = response.text.strip()
                            add_to_history_bot("", f"/img {prompt_response_text}")
                        generating = await channel.send("Generating image...")
                        add_to_history("System", f"Generating image: {prompt_response_text}")
                        if HUGGING_FACE_API == "HUGGING_FACE_API_KEY":
                            add_to_history("Error", "Failed to generate image! Invalid API Key.")
                            await channel.send("Failed to generate image! Invalid API Key.")
                            print("Failed to generate image! Invalid API Key, Please enter a valid hugging face API Key into system/config.py!")
                        else:
                            # Image Generation (Using Hugging Face Stable Diffusion)
                            api_key = HUGGING_FACE_API

                            url = f'https://api-inference.huggingface.co/models/{Image_Model}'
                            headers = {
                                'Authorization': f'Bearer {api_key}'
                            }
                            data = {
                                'inputs': prompt_response_text
                            }

                            response = requests.post(url, headers=headers, json=data)

                            if response.ok:
                                image_path = "system/RAM/gen-image/generated_image.png"
                                os.makedirs(os.path.dirname(image_path), exist_ok=True)
                                with open(image_path, 'wb') as f:
                                    f.write(response.content)
                                if add_watermark_to_generated_image:
                                    add_watermark("system/RAM/gen-image/generated_image.png", "system/RAM/gen-image/generated_image.png")
                                print("Image saved successfully as 'generated_image.png'!")

                                # Analyze the image
                                file_extension = image_path.split('.')[-1].lower()
                                if file_extension == 'jpg':
                                    file_extension = 'jpeg'
                                file_path = os.path.join('system/RAM/read-img', f'image.{file_extension}')

                                try:
                                    # Design and send the embed
                                    await generating.delete()
                                    embed = discord.Embed(title="Generated Image!",
                                                        description=f"{prompt_response_text}",
                                                        color=embed_colors)
                                    embed.set_image(url="attachment://generated_image.png")
                                    embed.set_footer(text=f"Generated by {NAME}")
                                    await channel.send(file=discord.File(image_path), embed=embed)
                                    add_to_history("Generated Image Details", response_text)
                                    await send_message(channel, response_text)

                                    os.remove(image_path)
                                except Exception as e:

                                    await channel.send("Error sending the generated image, Please try again later.")
                                    add_to_history("System", f"Error sending the generated image the image: {str(e)}")
                                    print(f"Error sending image: {e}")
                                    history = get_conversation_history(message) # Use the function to get history
                                    full_prompt = f"{history}\nError sending the image: {str(e)}"
                                    response = model.generate_content(full_prompt)  # Using the original language model
                                    response_text = response.text.strip()
                                    add_to_history_bot("", response_text)

                            else:
                                print('Error:', response.status_code, response.text)
                                add_to_history("Error", f"Failed to generate image: {response.status_code} | {response.text}")
                                await channel.send("An error occurred while generating the image.")

                    elif response_text.startswith("/object"):
                        # Object Detection (Using Hugging Face DETR)
                        API_URL = f"https://api-inference.huggingface.co/models/{Object_Detection_Model}"
                        headers = {"Authorization": f"Bearer {HUGGING_FACE_API}"}
                        response = requests.post(API_URL, headers=headers, data=image_files['data'])
                        try:
                            response.raise_for_status()
                            output = response.json()
                            print("Results:", output)
                            add_to_history("Object Detection Results", f"{output}")
                            # Create the 'system/RAM/annotate-img' directory if it doesn't exist
                            os.makedirs('system/RAM/object-img', exist_ok=True)
                            annotated_image_path = os.path.join('system/RAM/object-img', f"{os.path.basename(attachment.filename)}.object.jpg") 
                            # Create a figure and axes
                            fig, ax = plt.subplots(1)
                            ax.imshow(img)
                            # Draw bounding boxes
                            for prediction in output:
                                bbox = prediction['box']
                                score = prediction['score']
                                label = prediction['label']
                                rect = patches.Rectangle((bbox['xmin'], bbox['ymin']), 
                                                        bbox['xmax'] - bbox['xmin'], 
                                                        bbox['ymax'] - bbox['ymin'], 
                                                        linewidth=1, edgecolor='C'+str(len(output)), facecolor='none')
                                ax.add_patch(rect)
                                ax.text(bbox['xmin'], bbox['ymin']-10, label, color='C'+str(len(output)), fontsize=8)
                            plt.savefig(annotated_image_path)
                            plt.close()
                            response_obj_details = model_object.generate_content(f"{output}")  # Using the original language model
                            response_text_obj = response_obj_details.text.strip()
                            # Send the annotated image as a Discord Embed
                            embed = discord.Embed(title="Objects in the Image!",
                                description=f"{response_text_obj}",
                                color=embed_colors)
                            embed.set_image(url=f"attachment://{os.path.basename(annotated_image_path)}")
                            file = discord.File(annotated_image_path)
                            await channel.send(file=file, embed=embed)
                            # Clean up the annotated image file
                            os.remove(annotated_image_path)
                        except requests.exceptions.HTTPError as err:
                            print(f"Error: {err}")
                            await channel.send(f"An error occurred while detecting objects the image: {err}")
                    else:
                        await send_message(channel, response_text)

                if additional_details:

                    try:
                        response2 = model_pro.generate_content(image_file_objects)
                        response_text2 = response2.text.strip()
                        add_to_history("Additional Image details", response_text2)
                        print("Used model Pro")
                        print(" ")
                    except Exception as e3:
                        try:
                            response2 = model_V.generate_content(image_file_objects)
                            response_text2 = response2.text.strip()
                            add_to_history("Additional Image details", response_text2)
                            print("Used Model Pro Advanced")
                            print(" ")
                        except Exception as e2:
                            try:
                                response2 = model_V2.generate_content(image_file_objects)
                                response_text2 = response2.text.strip()
                                add_to_history("Additional Image details", response_text2)
                                print("Used Model Flash")
                                print(" ")
                            except Exception as e:
                                print(f"Failed to run Model Flash, Please try again Later | ERROR: {e}")
                                print(" ")
                            print(f"Failed running Model Pro Advanced, Running Model Flash | ERROR: {e2}")
                            print(" ")
                        print(f"Failed running Model Pro, Running Model Pro Advanced | ERROR: {e3}")
                        print(" ")
                break
            except Exception as e:
                unnecessary_error = await unnecessary_error(e)
                await debug_error(e, "Processing an Image", channel)
                if not unnecessary_error:
                    break

            finally:
                if response_text:
                    print(f"DEBUG: Final response text in 'finally' block: {response_text}")
                # Check if the directory exists
                if os.path.exists('system/RAM/read-img'):
                    # Iterate over each file in the directory and delete it
                    for filename in os.listdir('system/RAM/read-img'):
                        file_path = os.path.join('system/RAM/read-img', filename)
                        try:
                            if os.path.isfile(file_path) or os.path.islink(file_path):
                                os.unlink(file_path)  # Remove file or symbolic link
                            elif os.path.isdir(file_path):
                                shutil.rmtree(file_path)  # Remove directory and its contents
                        except Exception as e:
                            print(f"Error deleting {file_path}: {e}")
                else:
                    print("Directory 'system/RAM/read-img' does not exist.")
                break

# Function Definitions
def upload_to_gemini(path, mime_type=None):
    """Uploads the given file to Gemini."""
    file = genai.upload_file(path, mime_type=mime_type)
    print(f"Uploaded file '{file.display_name}' as: {file.uri}")
    return file

def wait_for_files_active(files):
    """Waits for the given files to be active."""
    print("Waiting for file processing...")
    for name in (file.name for file in files):
        file = genai.get_file(name)
        while file.state.name == "PROCESSING":
            print(".", end="", flush=True)
            time.sleep(10)
            file = genai.get_file(name)
        if file.state.name != "ACTIVE":
            raise Exception(f"File {file.name} failed to process")
    print("...all files ready")
    print()

async def handle_media_attachment(attachment, channel, prompt=None, message=None):
    """Handles the video and audio attachment processing and deletion."""
    file_extension = attachment.filename.split('.')[-1].lower()
    video_formats = ('mp4', 'avi', 'mkv', 'mov')

    # Determine whether it's a video or audio file
    is_video = file_extension in video_formats
    media_type = "video" if is_video else "audio"

    # Set the directory based on media type
    directory = 'system/RAM/read-vid' if is_video else 'system/RAM/read-audio'
    os.makedirs(directory, exist_ok=True)

    file_path = os.path.join(directory, 'media_file.' + file_extension)
    while True:
        try:
            media_data = await attachment.read()
            with open(file_path, 'wb') as file:
                file.write(media_data)

            # Upload the media to Gemini
            try:
                async with message.channel.typing():
                    gemini_file = upload_to_gemini(file_path, mime_type=f'{media_type}/{file_extension}')
                    wait_for_files_active([gemini_file])

                    # Prepare the text input (just the history)
                    display_name = member_custom_name.get(message.author.id, message.author.display_name)
                    member_name = display_name
                    history = get_conversation_history(message) # Use the function to get 
                    text_prompt = f"Conversation History Memory: |{history}|\n{member_name}: {prompt}"

                    # Generate response
                    response_text = ""  # Initialization
                    try:
                        # Send the message with the file URI and user prompt
                        if not gemini_api_key_verified:
                            await message.channel.send("Please set up the API key at [Google AI Studio](https://aistudio.google.com/apikey)")
                            await message.add_reaction('🔑')
                            return
                        response = chat_session.send_message([gemini_file, text_prompt])
                        response_text = response.text.strip()  # Assigning response_text here
                        add_to_history_bot("", response_text)
                        await send_message(channel, response_text)
                    except Exception as e:
                        print(f"Failed analyzing {media_type.capitalize()} | ERROR: {e}")

                if additional_details:
                    try:
                        response2 = model_vid.generate_content(gemini_file)
                        response_text2 = response2.text.strip()
                        add_to_history("Additional Media details", response_text2)
                        print("Used model Pro")
                    except Exception as e3:
                        try:
                            response2 = model_vid_a.generate_content(gemini_file)
                            response_text2 = response2.text.strip()
                            add_to_history("Additional Media details", response_text2)
                            print("Used Model Pro Advanced")
                        except Exception as e2:
                            try:
                                response2 = model_vid_flash.generate_content(gemini_file)
                                response_text2 = response2.text.strip()
                                add_to_history("Additional Media details", response_text2)
                                print("Used Model Flash")
                            except Exception as e:

                                print(f"Failed to run Model Flash, Please try again Later | ERROR: {e}")
                            print(f"Failed running Model Pro Advanced, Running Model Flash | ERROR: {e2}")
                        print(f"Failed running Model Pro, Running Model Pro Advanced | ERROR: {e3}")
                break

            except Exception as e:
                unnecessary_error = await unnecessary_error(e)
                if is_video:
                    await debug_error(e, "Processing a Video", channel)
                else:
                    await debug_error(e, "Processing an Audio", channel)
                if not unnecessary_error:
                    break

            finally:
                attempts = 0
                while True:
                    try:
                        # Delete the temporary media file
                        if os.path.exists(file_path):
                            os.remove(file_path)
                            print(f"Deleted file: {file_path}")
                        break
                    except PermissionError as e:
                        if "[WinError 32]" in str(e) or "used by another process" in str(e):
                            attempts += 1
                            print(f"Attempt {attempts}: File in use, retrying in 0.5 seconds...")
                            time.sleep(0.5)  # Short delay before retrying
                        else:
                            print(f"Error deleting file: {e}")
                            break

        except Exception as e:
            await channel.send(f"Error reading the attachment: {str(e)}")
            print(f"Error reading attachment: {e}")

async def handle_files_attachment(attachment, channel, prompt=None, message=None):
    """Handles various text-based file attachments and processes them."""

    file_extension = attachment.filename.split('.')[-1].lower()
    supported_file_formats = ('pdf', 'docx', 'md', 'py', 'js', 'bat', 'xlsx', 'pptx', 'csv', 'txt', 'json', 'log', 'html', 'css', 'mcmeta')

    if file_extension not in supported_file_formats:
        await channel.send("Unsupported file format. Please upload a supported file.")
        return
    directory = 'system/RAM/read-files'
    file_path = os.path.join(directory, attachment.filename)

    # Create the directory if it doesn't exist
    if not os.path.exists(directory):
        os.makedirs(directory)
    while True:
        try:
            # Save the file locally
            file_data = await attachment.read()
            with open(file_path, 'wb') as file:
                file.write(file_data)

            async with message.channel.typing():
                # Verify that the file was saved correctly
                if not os.path.exists(file_path):
                    raise FileNotFoundError(f"File {file_path} was not saved properly.")

                # Get file details
                file_size = os.path.getsize(file_path)
                creation_time = os.path.getctime(file_path)
                modification_time = os.path.getmtime(file_path)

                file_details = (
                    f"File: {attachment.filename}\n"
                    f"Size: {file_size} bytes\n"
                    f"Created: {datetime.datetime.fromtimestamp(creation_time)}\n"
                    f"Modified: {datetime.datetime.fromtimestamp(modification_time)}\n"
                )
                add_to_history("File Details", file_details)

                # Upload the file to Google Generative AI after saving it locally
                file = genai.upload_file(path=file_path)

                # Extract text content based on the file type
                text_content = ""

                if file_extension == 'docx':
                    # Extract text from DOCX files
                    try:
                        doc = docx.Document(file_path)
                        for paragraph in doc.paragraphs:
                            text_content += paragraph.text + "\n"
                    except Exception as e:

                        print(f"Error extracting text from .docx using docx module: {e}")
                        await channel.send(f"Error processing .docx file: Could not extract text.")
                        return
                elif file_extension == 'xlsx':
                    # Extract text from XLSX or XLS files
                    try:
                        workbook = openpyxl.load_workbook(file_path)
                        for sheet in workbook.sheetnames:
                            sheet_data = workbook[sheet]
                            for row in sheet_data.iter_rows(values_only=True):
                                text_content += ' '.join([str(cell) for cell in row]) + "\n"
                    except Exception as e:

                        print(f"Error extracting text from .xlsx: {e}")
                        await channel.send(f"Error processing .xlsx file: Could not extract text.")
                        return
                elif file_extension == 'pptx':
                    # Extract text from PPTX or PPT files
                    try:
                        presentation = pptx.Presentation(file_path)
                        for slide in presentation.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text_content += shape.text + "\n"
                    except Exception as e:

                        print(f"Error extracting text from .pptx or .ppt: {e}")
                        await channel.send(f"Error processing .pptx or .ppt file: Could not extract text.")
                        return
                elif file_extension == 'mcmeta':
                    # Read MCMETA files (usually small JSON-like files)
                    try:
                        with open(file_path, 'r') as mcmeta_file:
                            text_content = mcmeta_file.read()
                    except Exception as e:

                        print(f"Error reading .mcmeta file: {e}")
                        await channel.send(f"Error processing .mcmeta file: Could not read content.")
                        return

                add_to_history("System", f"{attachment.filename} received and processed")

                display_name = member_custom_name.get(message.author.id, message.author.display_name)
                member_name = display_name
                # Generate a response based on the file content and the prompt
                history = get_conversation_history(message)
                full_prompt = f"Conversation History Memory: |{history}|\n{member_name}'s {attachment.filename}{f': {prompt}' or ''}"

                if not gemini_api_key_verified:
                    await message.channel.send("Please set up the API key at [Google AI Studio](https://aistudio.google.com/apikey)")
                    await message.add_reaction('🔑')
                    return

                try:
                    if file_extension in ['docx', 'xlsx', 'xls', 'pptx', 'ppt', 'mcmeta']:
                        response = chat_session.send_message([f"{attachment.filename} DETAILS: '{text_content}' | ", full_prompt])
                    else:
                        response = chat_session.send_message([file, full_prompt])
                    response_text = response.text.strip()
                    if preview_code_output:
                        response_text = re.sub(r'``` *([a-zA-Z]+)', r'```\1', response_text)
                        response_text = re.sub(
                            r'```python\n(.*?)```\n+```*\n*(.*?)```',
                            r'```python\n\1```\nOutput:\n```bash\n\2```',
                            response_text,
                            flags=re.DOTALL
                        )

                except Exception as e:
                    print(f"Failed analyze file | ERROR: {e}")
                    response_text = "Sorry, I'm having trouble processing your request right now. Please try again later."
                # Add to conversation history for the main file details
                add_to_history_bot("", response_text)
                # Send the main response to the channel
                asyncio.create_task(send_message(channel, response_text))
            break

        except FileNotFoundError as fnf_error:
            await channel.send("Sorry, I'm having trouble processing your request right now. Please try again later.")
            print(f"FileNotFoundError: {fnf_error}")
            break
        except Exception as e:
            unnecessary_error = await unnecessary_error(e)
            await debug_error(e, "Processing a file", channel)
            if not unnecessary_error:
                break

        finally:
            # Remove the file from the system after processing
            if os.path.exists(file_path):
                os.remove(file_path)
                

img_path = f"system//RAM//read-img"  
os.makedirs(os.path.dirname(img_path), exist_ok=True)

skip_ffmpeg_check = False

async def generate_tts(message, response_text):
  global vc_voice, VOICES, ffmpeg_path, skip_ffmpeg_check
  print(f"Prompt before filtering: ({response_text})")

  if os.path.exists(ffmpeg_path) or skip_ffmpeg_check:
      TEXT = response_text
      TEXT = re.sub(r'<.*?>', '', TEXT)  #Remove < > tags only

      #Handle Emoji Only. (Instead of deleting other text)
      TEXT = re.sub(r'[^\w\s\.,!?"\':;-]', '', TEXT)

      TEXT = TEXT.strip()

      if isinstance(TEXT, bytes):
          TEXT = TEXT.decode('utf-8', errors='ignore')

      if not isinstance(TEXT, str):
          print(f"generate_tts ERROR: Text must be a string, not a {type(TEXT)} object. Check text contents before conversion")
          add_to_history("System", "Error during TTS generation. The generated text cannot be interpreted. Please contact an administrator if necessary.")
          return False
      print(f"Prompt after filtering: ({TEXT})")


      voice_index = vc_voice - 1
      if voice_index < 0 or voice_index >= len(VOICES):
          voice_index = 0

      VOICE = VOICES[voice_index]
      print(f"Chosen voice: {VOICE}, type: {type(VOICE)}")
      OUTPUT_FILE = "system/RAM/vc/Generated_voice.wav"

      try:
          print(f"Generating TTS for: {TEXT}")
          communicate = edge_tts.Communicate(TEXT, VOICE)
          print(f"communicate object: {communicate}")
          print("Saving TTS...")
          await communicate.save(OUTPUT_FILE)
          print("TTS file saved!")

          file_size = os.path.getsize(OUTPUT_FILE)
          print(f"Generated file size: {file_size} bytes")
          if file_size == 0:
              print("Warning: Generated file is empty.")
          return True

      except Exception as e:
          print(f"Error during TTS generation: {e}")
          return False

  else:
      print("Failed to generate VC. FFmpeg not found.")
      return False

async def play_tts(message, response_text):
  global vc_voice, VOICES, ffmpeg_path, skip_ffmpeg_check

  if os.path.exists(ffmpeg_path) or skip_ffmpeg_check:
      if not message.guild.voice_client:
          print("You need to be in a voice channel to play the TTS.")
      else:
          await generate_tts(message, response_text)
          OUTPUT_FILE = "system/RAM/vc/Generated_voice.wav"

          # Check if the author is in a voice channel
          if message.author.voice:
              voice_channel = message.author.voice.channel
              voice_client = message.guild.voice_client

              # Connect to the voice channel if not already connected
              if not voice_client:
                  voice_client = await voice_channel.connect()

              # Play the audio file if no audio is playing
              if not voice_client.is_playing():
                  voice_client.play(discord.FFmpegPCMAudio(OUTPUT_FILE, executable=ffmpeg_path))
              else:
                  print("Already playing audio, please wait.")
  else:
      print("Failed to generate VC. FFmpeg not found.")

ins2 = '''
    To help the user perform searches and analyze links using three main commands: "/search" for general web searches, "/search*yt" for YouTube video searches, and "/search [link]" for analyzing a provided link.

    ### Instructions for /search Command:
    1. Use "/search" for general web searches when the user does not specify a particular platform or type of content.
    2. The prompt after "/search" should contain the user's query. For example:
        - User: "Search for Google's Gemini API docs."
        - Response: "/search google gemini api docs"
    3. If the user sends a broad query, respond with the "/search" command and the query, ensuring clarity. 
    4. Only respond with the generated command in the format: "/search [query]". Do not explain or describe the process unless explicitly requested by the user.

    ### Instructions for /search*yt Command:
    1. Use "/search*yt" when the user requests a search for a YouTube video or implies they are looking for video content.
    2. The prompt after "/search*yt" should contain the search query. For example:
        - User: "Search YouTube for Python tutorials."
        - Response: "/search*yt Python tutorials"
    3. If the user specifies "search YouTube for...", or mentions a video-related term, respond with "/search*yt [query]".
    4. Do not include explanations; only return the search command.

    ### Instructions for /search [link] Command (Link Analysis):
    1. If the user provides a URL or link and requests you to analyze it, use the "/search [link]" command to process it.
    2. The URL must be included directly in the "/search" command. For example:
        - User: "Analyze this link: https://example.com/article."
        - Response: "/search https://example.com/article"
    3. Do not attempt to summarize, interpret, or explain the process of analyzing the link unless explicitly asked.
    4. Always ensure the link provided is included as-is in the "/search" command.
    5. If the user provides a link but doesn't explicitly ask for analysis, infer their intent based on the context and default to "/search [link]" if no other command fits.

    ### Important Notes:
    - **General Search vs. Link Analysis**:
        - Use "/search [query]" for general web queries.
        - Use "/search [link]" when the user provides a link or URL for analysis.
    - **YouTube Video Search**:
        - Use "/search*yt" for YouTube video-related searches.
    - **Clear and Concise Responses**:
        - Do not provide explanations or details about the commands or search process. Simply respond with the appropriate command and the user's query or link.

    ### Example Scenarios:
    1. **General Search**:
        - User: "Find the latest trends in machine learning."
        - AI Response: "/search latest trends in machine learning"

    2. **YouTube Video Search**:
        - User: "Search YouTube for funny dog videos."
        - AI Response: "/search*yt funny dog videos"

    3. **Link Analysis**:
        - User: "Analyze this link: https://example.com/AI-article."
        - AI Response: "/search https://example.com/AI-article"
        
    4. **Ambiguous Query**:
        - User: "Search for TensorFlow tutorials."
        - AI Response: "/search TensorFlow tutorials"
    
    5. **Ambiguous Query but for videos**:
        - User: "Search for TensorFlow video tutorials."
        - AI Response: "/search*yt TensorFlow tutorials"
        
    6. **Default Behavior**:
        - For unsupported scenarios or unclear instructions, default to interpreting the intent as a general search ("/search [query]") unless explicitly clarified by the user.

    Your task is to ensure accurate command generation based on the user's query, whether it's a general search, YouTube search, or link analysis. Follow the examples above for guidance, and do not include additional commentary unless explicitly requested by the user.
        ''',

logging.basicConfig(level=logging.DEBUG)

global_message = None

def clean_redundant_links(text):
  """
  Replaces Markdown links where the link text is the same as the URL
  with just the plain URL.

  Args:
      text: The input string containing Markdown links.

  Returns:
      The text with redundant Markdown links replaced by plain URLs.
  """
  link_pattern = r'\[(https?://[^\s)]+)\]\(\1\)'

  def replace_link(match):
    return match.group(1)

  return re.sub(link_pattern, replace_link, text)

def process_text(text):
    """
    Checks for the presence of a specific redundant link pattern and conditionally
    cleans the text.

    Args:
        text: The input string.

    Returns:
        The processed text, either cleaned or original.
    """

    pattern_to_check = r'\[(https?://[^\s)]+)\]\(\1\)'
    is_redundant_links =  bool(re.search(pattern_to_check, text))
    if is_redundant_links:
        return clean_redundant_links(text)
    else:
        return text


# Event handler to process messages
@bot.event
async def on_message(message):
    global global_message
    global_message = message
    if message.stickers or message.author.bot:
        return
    user = message.author.name
    user_settings = get_user_settings(user) # Call the function directly
    # Create a new task for processing the message
    asyncio.create_task(process_message(message, user_settings))
    # Don't wait for the task to finish; let it run concurrently
    return

async def process_message(message, user_settings):
    global model
    global vc_voice
    global VOICES
    global sys_security, gen_config, genai_model, ins, model_temperature, chat_session, ins2, ins, full_tutor_ins
    global ai_toggle_per_channel
    # Retrieve the custom name if set, otherwise use the default display name
    display_name = member_custom_name.get(message.author.id, message.author.display_name)
    channel_id = message.channel.id
    add_to_history(display_name, message.content)
    if smart_recognition and ai_toggle_per_channel.get(channel_id, False) and not message.content.startswith("/") and message.guild:
        if dev_DEBUG:
            print("Smart recognition is enabled.")
        smart_recognition_model = genai.GenerativeModel(
            model_name="gemini-1.5-flash-latest",
            generation_config=gen_config,
            system_instruction=f"""
                You are a smart recognition engine designed to determine if the user is interacting with the bot or another person. The bot’s name is {NAME}. You have access to the entire conversation history, and based on this, you need to decide whether the user is talking to the bot or someone else. 

                **Your task is simple**:
                - Respond with `1` if the user is talking to the bot.
                - Respond with `0` if the user is talking to someone else.

                **Important guidelines**:
                - Only respond with `1` or `0` — no other responses.
                - You can see the conversation history, so use it carefully to determine if the user is talking to the bot or a friend.
                - Pay special attention if the user asks the bot not to respond until a certain prompt (e.g., "don’t answer my first and second prompt but answer my third prompt"). In this case, return `0` for the first and second prompts and `1` for the third prompt.
                - If the user was previously talking to a friend but then asks the bot to answer a specific question, you should recognize this and let the bot answer when appropriate.
                - A user might not always mention the bot's name for it to respond, so be alert to conversational cues.
                - If a user repeats a greeting or a question and no one responds, it’s likely they are trying to get the bot’s attention.
                - If the user wasn’t talking to a friend and is engaged in a conversation, it’s likely with the bot.

                **Tip**: Pay attention to shifts in the conversation. If the user hasn’t been talking to a friend and suddenly begins talking, it’s a strong indication that the user is talking to the bot.

                **Be cautious and listen carefully** to the conversation history to make the best judgment. Your ability to detect whether the user is talking to the bot or a person will depend on your careful analysis of the chat history. Good luck!
            """,
            safety_settings=sys_security
        )
        if dev_DEBUG:
            print("Smart recognition model initialized.")
        smart_recognition_chat = smart_recognition_model.start_chat(history=[])
        if dev_DEBUG:
            print("Smart recognition chat started.")

        async def smart_recognition_check(message):
            if dev_DEBUG:
                print(f"Checking smart recognition for message: {message.content}")
            channel_name = message.channel.name
            channel_id = message.channel.id
            history = get_conversation_history(message)  # Use the function to get history
            if dev_DEBUG:
                print(f"Conversation history: {history}")
            full_prompt = f"From the chat named '{channel_name}|ID:{channel_id}': \n\n{history}"
            if dev_DEBUG:
                print(f"Full prompt: {full_prompt}")

            response = smart_recognition_chat.send_message(full_prompt)
            if dev_DEBUG:
                print(f"Raw response from model: {response.text.strip()}")
            raw_results = response.text.strip()
            if "1" in raw_results:
                results = 1
                if dev_DEBUG:
                    print("User is talking to the bot.")
            elif "0" in raw_results:
                results = 0
                if dev_DEBUG:
                    print("User is talking to someone else.")
            else:
                results = None
                if dev_DEBUG:
                    print(f"Unexpected result from smart recognition. | {raw_results}")

            return results

        result = await smart_recognition_check(message)
        if dev_DEBUG:
            print(f"Smart recognition result: {result}")

        if result == 1:
            if dev_DEBUG:
                print("Proceeding with bot interaction.")
            pass
        elif result == 0:
            if dev_DEBUG:
                print("Returning, user is not talking to the bot.")
            return

    # Check if AI is toggled on for this channel
    if ai_toggle_per_channel.get(channel_id, False) and not message.content.startswith("/") and message.guild:

        if not gemini_api_key_verified:
            await message.channel.send("Please set up the API key at [Google AI Studio](https://aistudio.google.com/apikey)")
            print("Please set up the API key at https://aistudio.google.com/apikey")
            await message.add_reaction('🔑')
            return
        elif is_youtube_url(message.content):
            async with message.channel.typing():
                transcript = await handle_youtube_url(message.content, message.channel, prompt=message.content)
                add_to_history("System", f"Analyzed the YouTube URL and retrieved transcript: {transcript}")
                history = get_conversation_history(message)
                display_name = message.author.display_name
                full_prompt = f"{history}"
                response = chat_session.send_message(full_prompt)
                if Model_Debug:
                    print(f"Raw response from model: {response}")
                    print(f"Model configuration: {model}")
                response_text = response.text.strip()
                if preview_code_output:
                    response_text = re.sub(r'``` *([a-zA-Z]+)', r'```\1', response_text)
                    response_text = re.sub(
                        r'```python\n(.*?)```\n+```*\n*(.*?)```',
                        r'```python\n\1```\nOutput:\n```bash\n\2```',
                        response_text,
                        flags=re.DOTALL
                    )
                response_text = process_text(response_text)
                add_to_history_bot("", response_text)
                if sync_voice_with_text:
                    if tts_toggle and not vc_AI:  # If TTS is enabled
                        if response_text:
                            asyncio.create_task(play_tts(message, response_text))
                        asyncio.create_task(send_message(message.channel, f"- <:YT:1312077123971842080> **Successfully Analyzed Youtube URL:**\n{response_text}"))
                        return
                    else:
                        asyncio.create_task(send_message(message.channel, f"- <:YT:1312077123971842080> **Successfully Analyzed Youtube URL:**\n{response_text}"))
                        return
                else:
                    asyncio.create_task(send_message(message.channel, f"- <:YT:1312077123971842080> **Successfully Analyzed Youtube URL:**\n{response_text}"))
                    if tts_toggle and not vc_AI:  # If TTS is enabled
                        asyncio.create_task(play_tts(message, response_text))
                        return
                    return

        elif message.attachments:
            for attachment in message.attachments:
                valid_attachments = [
                attachment for attachment in message.attachments
                if attachment.filename.lower().endswith(('png', 'jpg', 'jpeg', 'gif'))
            ]
            if valid_attachments:
                asyncio.create_task(handle_image_attachment(valid_attachments, message.channel, prompt=message.content, message=message))
                return
            elif attachment.filename.lower().endswith(('mp4', 'avi', 'mkv', 'mov', 'mp3', 'wav', 'aac')):
                asyncio.create_task(handle_media_attachment(attachment, message.channel, message=message))
                return
            elif attachment.filename.lower().endswith(('pdf', 'docx', 'md', 'py', 'js', 'bat', 'xlsx', 'pptx', 'csv', 'txt', 'json', 'log', 'html', 'css', 'mcmeta')):
                asyncio.create_task(handle_files_attachment(attachment, message.channel, prompt=message.content, message=message))
                return
            else:
                print(f"Unsupported file format: {attachment.filename} | {attachment.filename.lower()} | From: {message.author.display_name} | ID: {message.author.id}")
                asyncio.create_task(send_message(message.channel, "Unsupported file format. Please upload a supported file format."))
                return

        else:
            def extract_check_info(response_text):
                """
                Extracts whether the prompt is for image generation and the prompt text itself.
                
                Args:
                - response_text (str): The response from the LLM model.

                Returns:
                - bool: True if the prompt is for image generation, False otherwise.
                - str: The extracted prompt (if valid), or an empty string if not.
                """
                # Regular expression to match the pattern "1|<prompt>" or "0"
                match = re.match(r"(\d)\|(.+)", response_text.strip())  # Match 1|<prompt> or 0
                
                if match:
                    # Extract the 1 or 0 and convert to True or False
                    result = match.group(1) == '1'
                    # Extract the prompt after the "/img" command
                    prompt = match.group(2).strip()
                    return result, prompt
                elif response_text.strip() == '0':
                    return False, ''
                else:
                    # In case the response does not follow the expected pattern
                    return None, ''
            if_image_generation_prompt_model = genai.GenerativeModel(
                model_name="gemini-1.5-flash",
                generation_config=gen_config,
                system_instruction='Your porpuse is to check if the prompt is an image generation prompt, usually the prompt must use the command `/img` to generate an image, but dont get tricked by other commands and if its just explaining the command or not, if the prompt contains the command `/img` and the user wants to generate an image, respond with "1" if not, respond with "0" but also beside it put the prompt that the user wants to generate like "/img a cute fox on an island" and you will respond with "1|a cute fox on an island", and dont say anything than "1|[prompt]" and "0"',
                safety_settings=sys_security
            )
            if_music_generation_prompt_model = genai.GenerativeModel(
                model_name="gemini-1.5-flash",
                generation_config=gen_config,
                system_instruction='''
                    Your purpose is to check if the prompt is a music generation prompt. Usually, the prompt must contain the command `/music` to generate or play music. Do not get tricked by other commands or if it's just explaining the command, but if the prompt contains the command `/music` and the user wants to generate or play music, respond with "1". If not, respond with "0", but also beside it include the prompt that the user wants to generate or play, like `/music a relaxing song` and you will respond with "1|a relaxing song". 
                    Only respond with "1|[prompt]" or "0", nothing else.
                ''',
                safety_settings=sys_security
            )
            if_youtube_search_prompt_model = genai.GenerativeModel(
                model_name="gemini-1.5-flash",
                generation_config=gen_config,
                system_instruction='''
                    Your purpose is to check if the prompt is asking to search for a YouTube video. Usually, the prompt will contain "/search*yt" or a variation of the command to request a YouTube search. 
                    If the prompt contains "/search*yt" and the user wants to search for a YouTube video, respond with "1" and the query they want to search for, like "/search*yt funny cat videos", you will respond with "1|funny cat videos".
                    If the prompt does not include "/search*yt" or is not asking for a YouTube search, respond with "0". 
                    Only respond with "1|[search query]" or "0", nothing else.
                ''',
                safety_settings=sys_security
            )

            if_general_search_prompt_model = genai.GenerativeModel(
                model_name="gemini-1.5-flash",
                generation_config=gen_config,
                system_instruction='''
                    Your purpose is to check if the prompt is a general search query (e.g., a web search).
                    The prompt will typically include "/search" followed by the user's search query.
                    If the prompt contains "/search" and the user wants to perform a general search, respond with "1" and the search query, like "/search how to make pizza", you will respond with "1|how to make pizza".
                    If the prompt does not include "/search" or is not asking for a general search, respond with "0". 
                    Only respond with "1|[search query]" or "0", nothing else.
                    and if the prompt was trying to search using code, try knowing what the prompt was trying to search and do "1|[search query]"
                ''',
                safety_settings=sys_security
            )



            while True:  # Loop to handle retries
                try:
                    async with message.channel.typing():
                        selected_model = user_settings['model']
                        selected_model_name = user_settings['model_name']
                        full_ins = f"{ins}\n\nWe have found the model your are running on! You are currently running on ({selected_model_name}), so if someone asks you what model you are running, tell them that its {selected_model_name}"
                        full_tutor_ins = f"{tutor_ins}\n\nWe have found the model your are running on! You are currently running on ({selected_model_name}), so if someone asks you what model you are running, tell them that its {selected_model_name}"
                        models = {
                            'models/learnlm-1.5-pro-experimental': {
                                'model_name': 'learnlm-1.5-pro-experimental',
                                'system_instruction': full_tutor_ins
                            },
                            'models/gemini-1.5-flash': {
                                'model_name': 'gemini-1.5-flash',
                                'system_instruction': full_ins,
                            },
                            'models/gemini-1.5-flash-002': {
                                'model_name': 'gemini-1.5-flash-002',
                                'system_instruction': full_ins,
                            },
                            'models/gemini-1.5-flash-8b': {
                                'model_name': 'gemini-1.5-flash-8b',
                                'system_instruction': full_ins,
                            },
                            'models/gemini-1.5-pro': {
                                'model_name': 'gemini-1.5-pro',
                                'system_instruction': full_ins,
                            },
                            'models/gemini-1.5-pro-002': {
                                'model_name': 'gemini-1.5-pro-002',
                                'system_instruction': full_ins,
                            },
                            'models/gemini-1.5-pro-latest': {
                                'model_name': 'gemini-1.5-pro-latest',  # Or the actual latest model name if different
                                'system_instruction': full_ins,
                            },
                            'models/gemini-1.5-flash-latest': {
                                'model_name': 'gemini-1.5-flash-latest', # Or the actual name
                                'system_instruction': full_ins,
                            },
                            'models/gemini-exp-1114': {
                                'model_name': 'gemini-exp-1114',
                                'system_instruction': full_ins,
                            },
                            'models/gemini-exp-1121': {
                                'model_name': 'gemini-exp-1121',
                                'system_instruction': full_ins,
                            },
                            'models/gemini-exp-1206': {
                                'model_name': 'gemini-exp-1206',
                                'system_instruction': full_ins,
                            },
                        }

                        selected_model = f"models/{selected_model}"

                        if selected_model in models:
                            # Debugging: Print selected model and current model details
                            if dev_DEBUG:
                                print(f"DEBUG: Selected model: {selected_model}")
                            if hasattr(model, 'model_name') and dev_DEBUG:
                                print(f"DEBUG: Current model name: {model.model_name}")
                            else:
                                if dev_DEBUG:
                                    print("DEBUG: Current model does not have a 'model_name' attribute. Model re-creation required.")

                            # Check if the current model matches the user's selected model
                            if not hasattr(model, 'model_name') or model.model_name != selected_model:
                                if dev_DEBUG:
                                    print("DEBUG: Model mismatch detected or model is not initialized. Re-creating the model...")
                                try:
                                    model_config = models[selected_model]
                                    if dev_DEBUG:
                                        print(f"DEBUG: Model configuration retrieved for {selected_model}: {model_config}")
                                    model = genai.GenerativeModel(
                                        model_name=model_config['model_name'],
                                        generation_config=gen_config,
                                        system_instruction=model_config['system_instruction'],
                                        safety_settings=sys_security,
                                        tools='code_execution' if preview_code_output else None
                                    )
                                    if dev_DEBUG:
                                        print(f"DEBUG: Model '{selected_model}' has been successfully created.")
                                    chat_session = model.start_chat()  # Start chat session *after* creating the model
                                    if dev_DEBUG:
                                        print("DEBUG: Chat session started.")
                                except Exception as e:
                                    if dev_DEBUG:
                                        print(f"ERROR: Failed to re-create the model. Exception: {e}")
                            else:
                                if dev_DEBUG:
                                    print("DEBUG: Selected model matches the current model. No re-creation needed.")
                        else:
                            print(f"ERROR: Invalid model selected: {selected_model}")
                            await message.channel.send(f"Error: Invalid model selected: {selected_model}")
                            return
                        history = get_conversation_history(message)
                        display_name = message.author.display_name
                        full_prompt = f"{history}"
                        response = chat_session.send_message(full_prompt)
 
                        if Model_Debug:
                            print(f"Raw response from model: {response}")
                            print(f"Model configuration: {model}")
                        response_text = response.text.strip()
                        if preview_code_output:
                            response_text = re.sub(r'``` *([a-zA-Z]+)', r'```\1', response_text)
                            response_text = re.sub(
                                r'```python\n(.*?)```\n+```*\n*(.*?)```',
                                r'```python\n\1```\nOutput:\n```bash\n\2```',
                                response_text,
                                flags=re.DOTALL
                            )
                        
                        if "/img" in response_text:
                            check_if_img = if_image_generation_prompt_model.generate_content(response_text)
                            is_image_gen, checked_img_prompt = extract_check_info(check_if_img.text)
                            prompt_response_text = ""
                            if is_image_gen:
                                check_prompt_response_text = checked_img_prompt
                            else:
                                if show_tokens:
                                    tokens = model.count_tokens(f"{response_text}")
                                    print(f"----------------------------------------")
                                    print(f"{Style.BRIGHT, Fore.WHITE}{tokens}{Style.RESET_ALL}")
                                    print(f"----------------------------------------")

                                add_to_history_bot("", response_text)
                                if tts_toggle and not vc_AI:
                                    # Check if the bot is in a voice channel
                                    if message.guild.voice_client is not None:  
                                        if sync_voice_with_text:
                                            if response_text:
                                                asyncio.create_task(play_tts(message, response_text))
                                            asyncio.create_task(send_message(message.channel, response_text))
                                            return
                                        else:
                                            asyncio.create_task(send_message(message.channel, response_text))
                                            asyncio.create_task(play_tts(message, response_text))
                                            return
                                    else:
                                        asyncio.create_task(send_message(message.channel, response_text))
                                        return
                                else:
                                    asyncio.create_task(send_message(message.channel, response_text))
                                    return
                            await message.channel.send("Generating image...")
                            add_to_history("System", f"Generating image: {check_prompt_response_text}")

                            if HUGGING_FACE_API == "HUGGING_FACE_API_KEY":
                                add_to_history("Error", "Failed to generate image! Invalid API Key.")
                                await message.channel.send("Failed to generate image! Invalid API Key.")
                                print("Failed to generate image! Invalid API Key, Please enter a valid hugging face API Key into system/config.py!")
                            else:
                                try:
                                    model_check_prompt_flash = genai.GenerativeModel(
                                        model_name="gemini-1.5-flash-latest",
                                        generation_config=gen_config,
                                        system_instruction="your porpuse is to only give out numbers to the prompts: if the prompt is a diffrent language than english only say `1` and if its an inappropriate prompt only say `2` and if its corrupted only say `3` and only say `4` if its good to go and ONLY say `5` if its REALLY inappropriate or smth and needs to be sent to the moderators to temp ban or timeout the user.",
                                        safety_settings=sys_security,
                                    )
                                    try:
                                        response_check = model_check_prompt_flash.generate_content(
                                            check_prompt_response_text
                                        )
                                        response_check_text_check = response_check.text.strip()
                                        print("Used Model Flash on check")
                                    except Exception as e:
                                        print(f"Failed to run image generation prompt check. | ERROR: {e}")
                                        await message.channel.send(
                                            "Failed to run image generation prompt check. Please try again later."
                                        )
                                        return
                                    if "1" in response_check_text_check:
                                        model_translation_gen_prompt_flash = genai.GenerativeModel(
                                            model_name="gemini-1.5-flash-latest",
                                            generation_config=gen_config,
                                            system_instruction="your porpuse is to only to translate any user's prompt language to english, and nothing else, you must not say anything unless its the translated prompt, just like google translate!",
                                            safety_settings=sys_security,
                                        )
                                        try:
                                            response_translation_gen_image = (
                                                model_translation_gen_prompt_flash.generate_content(
                                                    check_prompt_response_text
                                                )
                                            )
                                            (
                                                response_translation_gen_image_text_translate
                                            ) = response_translation_gen_image.text.strip()
                                            prompt_response_text = (
                                                f"{response_translation_gen_image_text_translate}"
                                            )
                                            print("Used Model Flash")
                                        except Exception as e:
                                            print(f"Failed to run translating image generation prompt. | ERROR: {e}")
                                            await message.channel.send("Sorry, Failed to translate prompt to english, please try an english prompt")
                                            return
                                        print("Translated prompt!")
                                    elif "2" in response_check_text_check or "5" in response_check_text_check:
                                        error_message = "I'm unable to create an image based on your request. Please make sure your prompt aligns with our image generation guidelines, which focus on creating safe and appropriate visuals."
                                        add_to_history("Error", error_message)
                                        await message.channel.send(error_message)
                                        return
                                    elif "3" in response_check_text_check:
                                        print("Corrupted image generation prompt")
                                        add_to_history(
                                            "Error",
                                            "Oops, something seems off with that prompt. Please try rephrasing it or using different keywords. I'm here to help if you need suggestions!",
                                        )
                                        await message.channel.send(
                                            "Oops, something seems off with that prompt. Please try rephrasing it or using different keywords. I'm here to help if you need suggestions!"
                                        )
                                        return
                                    elif "4" in response_check_text_check:
                                        print("Image generation prompt is safe!")
                                        prompt_response_text = check_prompt_response_text
                                    else:
                                        print(f"Checking prompt Error: {e}")
                                        add_to_history(
                                            "System",
                                            "Oops! Looks like the image generator took a coffee break ☕. Please try again in a moment!",
                                        )
                                        await message.channel.send(
                                            "Oops! Looks like the image generator took a coffee break ☕. Please try again in a moment!"
                                        )
                                        return
                                except Exception as e:
                                    print(f"Checking prompt Error: {e}")
                                    add_to_history(
                                        "System",
                                        "Oops! Looks like the image generator took a coffee break ☕. Please try again in a moment!",
                                    )
                                    await message.channel.send(
                                        "Oops! Looks like the image generator took a coffee break ☕. Please try again in a moment!"
                                    )
                                    return
                                # Image Generation (Using Hugging Face Stable Diffusion)
                                api_key = HUGGING_FACE_API

                                url = f"https://api-inference.huggingface.co/models/{Image_Model}"
                                headers = {"Authorization": f"Bearer {api_key}"}
                                data = {"inputs": prompt_response_text}

                                response = requests.post(url, headers=headers, json=data)

                                if response.ok:
                                    image_path = "system/RAM/gen-image/generated_image.png"
                                    os.makedirs(os.path.dirname(image_path), exist_ok=True)
                                    with open(image_path, "wb") as f:
                                        f.write(response.content)
                                    if add_watermark_to_generated_image:
                                        add_watermark("system/RAM/gen-image/generated_image.png", "system/RAM/gen-image/generated_image.png")
                                    print("Image saved successfully as 'generated_image.png'!")

                                    # Analyze the image
                                    file_extension = image_path.split(".")[-1].lower()
                                    if file_extension == "jpg":
                                        file_extension = "jpeg"

                                    try:
                                        # Design and send the embed
                                        embed = discord.Embed(
                                            title="Generated Image!",
                                            description=f"{prompt_response_text}",
                                            color=embed_colors,
                                        )
                                        embed.set_image(url="attachment://generated_image.png")
                                        embed.set_footer(text=f"Generated by {NAME}")
                                        await message.channel.send(
                                            file=discord.File(image_path), embed=embed
                                        )
                                        add_to_history("Generated Image Details", response_text)

                                        os.remove(image_path)
                                    except Exception as e:
                                        await message.channel.send(
                                            "Error sending the generated image, Please try again later."
                                        )
                                        add_to_history(
                                            "System",
                                            f"Error sending the generated image the image: {str(e)}",
                                        )
                                        print(f"Error sending image: {e}")
                                        history = get_conversation_history(
                                            message
                                        )  # Use the function to get history
                                        full_prompt = (
                                            f"{history}\nError sending the image: {str(e)}"
                                        )
                                        response = model.generate_content(
                                            full_prompt
                                        )  # Using the original language model
                                        response_text = response.text.strip()
                                        add_to_history_bot("", response_text)

                                else:
                                    print("Error:", response.status_code, response.text)
                                    add_to_history(
                                        "Error",
                                        f"Failed to generate image: {response.status_code} | {response.text}",
                                    )
                                    await message.channel.send(
                                        "An error occurred while generating the image."
                                    )
                        elif "/music" in response_text:
                            check_if_music = if_music_generation_prompt_model.generate_content(response_text)
                            is_music_gen, checked_music_prompt = extract_check_info(check_if_music.text)
                            prompt_response_text = ""
                            
                            if is_music_gen:
                                # If it's a valid music prompt, use the checked prompt
                                prompt = checked_music_prompt
                            else:
                                # If it's not a valid music prompt, process normally
                                if show_tokens:
                                    tokens = model.count_tokens(f"{response_text}")
                                    print(f"----------------------------------------")
                                    print(f"{Style.BRIGHT, Fore.WHITE}{tokens}{Style.RESET_ALL}")
                                    print(f"----------------------------------------")

                                add_to_history_bot("", response_text)
                                if tts_toggle and not vc_AI:
                                    # Check if the bot is in a voice channel
                                    if message.guild.voice_client is not None:  
                                        if sync_voice_with_text:
                                            if response_text:
                                                asyncio.create_task(play_tts(message, response_text))
                                            asyncio.create_task(send_message(message.channel, response_text))
                                            return
                                        else:
                                            asyncio.create_task(send_message(message.channel, response_text))
                                            asyncio.create_task(play_tts(message, response_text))
                                            return
                                    else:
                                        asyncio.create_task(send_message(message.channel, response_text))
                                        return
                                else:
                                    asyncio.create_task(send_message(message.channel, response_text))
                                    return


                            # Music Generation Logic (Integrated from generate_music function)
                            if HUGGING_FACE_API == "YOUR_HUGGING_FACE_API_KEY":
                                await message.channel.send(
                                    "Sorry, You have entered an Invalid Hugging Face API Key!"
                                )
                                return

                            await message.channel.send(
                                "Generating Music..."
                            )  # Indicate that music generation is starting

                            api_key = HUGGING_FACE_API
                            max_retries = 10
                            backoff_factor = 3

                            music_model = DEFAULT_MUSIC_MODEL  # Default model
                            url = f"https://api-inference.huggingface.co/models/{music_model}"
                            headers = {"Authorization": f"Bearer {api_key}"}
                            data = {"inputs": prompt}

                            def save_audio(response):
                                audio_dir = "system/RAM/gen-music"
                                os.makedirs(audio_dir, exist_ok=True)
                                audio_path = os.path.join(audio_dir, "generated_music.wav")
                                with open(audio_path, "wb") as f:
                                    f.write(response.content)
                                logging.info(
                                    f"Audio saved successfully as '{audio_path}'!"
                                )

                            def handle_error(response):
                                error_message = response.json().get(
                                    "error", "No error message"
                                )
                                if response.status_code == 503:
                                    logging.error(
                                        f"Service unavailable. Error: {error_message}"
                                    )
                                elif response.status_code == 429:
                                    logging.error(
                                        f"Rate limit exceeded. Error: {error_message}"
                                    )
                                else:
                                    logging.error(
                                        f"Failed to save audio. Status code: {response.status_code}, Error: {error_message}"
                                    )

                            def fetch_audio_with_retries(url, headers, data):
                                for attempt in range(max_retries):
                                    response = requests.post(
                                        url, headers=headers, json=data
                                    )
                                    if response.ok:
                                        save_audio(response)
                                        return True
                                    else:
                                        handle_error(response)
                                        if response.status_code in [503, 429]:
                                            wait_time = backoff_factor**attempt
                                            logging.info(
                                                f"Retrying in {wait_time} seconds..."
                                            )
                                            time.sleep(wait_time)
                                        else:
                                            break
                                logging.error(
                                    "Exceeded maximum retries or encountered a non-retryable error."
                                )
                                return False

                            success = fetch_audio_with_retries(url, headers, data)
                            if success:
                                audio_path = "system/RAM/gen-music/generated_music.wav"
                                file = discord.File(
                                    audio_path, filename="generated_music.wav"
                                )
                                await message.channel.send(file=file)
                                os.remove(audio_path)
                            else:
                                await message.channel.send(
                                    "An error occurred while generating the music. Please try again later."
                                )
                        elif "/search*yt" in response_text:  # For YouTube search
                            check_if_youtube = if_youtube_search_prompt_model.generate_content(response_text)
                            is_youtube_search, checked_youtube_prompt = extract_check_info(check_if_youtube.text)
                            
                            if is_youtube_search:
                                search_query = checked_youtube_prompt  # Valid YouTube search prompt
                            else:
                                if show_tokens:
                                    tokens = model.count_tokens(f"{response_text}")
                                    print(f"----------------------------------------")
                                    print(f"{Style.BRIGHT, Fore.WHITE}{tokens}{Style.RESET_ALL}")
                                    print(f"----------------------------------------")

                                add_to_history_bot("", response_text)
                                if tts_toggle and not vc_AI:
                                    # Check if the bot is in a voice channel
                                    if message.guild.voice_client is not None:  
                                        if sync_voice_with_text:
                                            if response_text:
                                                asyncio.create_task(play_tts(message, response_text))
                                            asyncio.create_task(send_message(message.channel, response_text))
                                            return
                                        else:
                                            asyncio.create_task(send_message(message.channel, response_text))
                                            asyncio.create_task(play_tts(message, response_text))
                                            return
                                    else:
                                        asyncio.create_task(send_message(message.channel, response_text))
                                        return
                                else:
                                    asyncio.create_task(send_message(message.channel, response_text))
                                    return
                            if web_search:
                                if google_search_api_verified:
                                    try:
                                        global search_google
                                        add_to_history_bot("", f"/search*yt {search_query}")
                                        search_results = search_google(search_query, site='https://www.youtube.com/', num_results=20)

                                        try:
                                            if dev_DEBUG:
                                                print(f"Youtube Search Results: {search_results}")
                                            add_to_history("Searched Youtube Videos", search_results)
                                            add_to_history("web_search_results", "Ok, tell the user about those YouTube videos")
                                            history = get_conversation_history(message) # Use the function to get history
                                            full_prompt = f"{history}\nSearch: {search_results}"
                                            response = chat_session.send_message(full_prompt)
                                            response_text = response.text.strip()
                                            response_text = process_text(response_text)
                                            add_to_history("Me", response_text)
                                            if not "Error Searching the web." in search_results:
                                                await send_message(message.channel, f"- <:YT:1312077123971842080> **Successfully Searched Youtube:**\n{response_text}")
                                            else:
                                                await send_message(message.channel, f"- 🚫 **Failed Searching Youtube:**\n{response_text}")
                                        except Exception as e:
                                            print(f"Error generating content: {e}")
                                            return "An error occurred while generating the response."

                                    except Exception as e:
                                        await message.channel.send("Error occurred during search Youtube. Try again later.")
                                        print(e)
                                else:
                                    await message.channel.send("Error: `Google Search API key` is not valid.")
                                    print("Error: Google Search API key is not valid.`")
                            else:
                                await message.channel.send("Error: Web Search is disabled, Cannot search Youtube...")
                                print("Error: Web Search is disabled, Cannot search Youtube...")
                        elif "/search" in response_text:  # For general search
                            check_if_general_search = if_general_search_prompt_model.generate_content(response_text)
                            is_general_search, checked_general_prompt = extract_check_info(check_if_general_search.text)
                            prompt_response_text = ""
                            
                            if is_general_search:
                                search_query = checked_general_prompt  # Valid general search prompt
                            else:
                                if show_tokens:
                                    tokens = model.count_tokens(f"{response_text}")
                                    print(f"----------------------------------------")
                                    print(f"{Style.BRIGHT, Fore.WHITE}{tokens}{Style.RESET_ALL}")
                                    print(f"----------------------------------------")

                                add_to_history_bot("", response_text)
                                if tts_toggle and not vc_AI:
                                    # Check if the bot is in a voice channel
                                    if message.guild.voice_client is not None:  
                                        if sync_voice_with_text:
                                            if response_text:
                                                asyncio.create_task(play_tts(message, response_text))
                                            asyncio.create_task(send_message(message.channel, response_text))
                                            return
                                        else:
                                            asyncio.create_task(send_message(message.channel, response_text))
                                            asyncio.create_task(play_tts(message, response_text))
                                            return
                                    else:
                                        asyncio.create_task(send_message(message.channel, response_text))
                                        return
                                else:
                                    asyncio.create_task(send_message(message.channel, response_text))
                                    return
                            if web_search:
                                if google_search_api_verified:
                                    try:
                                        add_to_history_bot("", f"/search {search_query}")
                                        if "http://" in search_query or "https://" in search_query:
                                            # Extract the link from the search query
                                            match = re.search(r'(https?://[^\s]+)', search_query)
                                            if match:
                                                site = match.group(0)  # Extracted link
                                                search_query = search_query.replace(site, "").strip()  # Remove the link from search_query
                                                web_search_results = search_google(search_query, site=site)
                                            else:
                                                # Fallback in case no valid link is found
                                                web_search_results = search_google(search_query)
                                        else:
                                            web_search_results = search_google(search_query)
                                        try:
                                            if dev_DEBUG:
                                                print(f"Web Search Results: {web_search_results}")
                                            add_to_history("Searched Web Results", web_search_results)
                                            add_to_history("Web Search Engine", "Ok, tell the user about that search")
                                            history = get_conversation_history(message) # Use the function to get history
                                            full_prompt = f"{history}\nSearch: {web_search_results}"
                                            response = chat_session.send_message(full_prompt)
                                            response_text = response.text.strip()
                                            response_text = process_text(response_text)
                                            add_to_history_bot("", response_text)
                                            if not "Error Searching the web." in web_search_results:
                                                await send_message(message.channel, f"- :mag: **Successfully Searched the web:**\n{response_text}")
                                            else:
                                                await send_message(message.channel, f"- 🚫 **Failed Searching the web:**\n{response_text}")
                                        except Exception as e:

                                            print(f"Error generating content: {e}")
                                            return "An error occurred while generating the response."
                                    except Exception as e:
                                        if e.response.status_code == 429:
                                            print(f"Rate limit hit! \nERROR: {e}")
                                            await send_message(message.channel, f"Sorry, it seems like you have reached the limit for the web search. Please try again later.")
                                            add_to_history("Failed-Search", f"An error occurred during search. Please try again later.\nSorry, it seems like you have reached the limit for the web search. Please try again later.\nERROR: {e}")
                                            return
                                        print(e)
                                        add_to_history("Failed-Search", f"An error occurred during search. Please try again later.\nERROR: {str(e)}")
                                else:
                                        await message.channel.send("Error: `Google Search API key` is not valid.")
                                        print("Error: Google Search API key is not valid.`")
                            else:
                                await message.channel.send("Error: Web Search is disabled, Cannot search the Web...")
                                print("Error: Web Search is disabled, Cannot search the Web...")
                                
                        else:
                            if show_tokens:
                                tokens = model.count_tokens(f"{response_text}")
                                print(f"----------------------------------------")
                                print(f"{Style.BRIGHT, Fore.WHITE}{tokens}{Style.RESET_ALL}")
                                print(f"----------------------------------------")

                            add_to_history_bot("", response_text)
                            if tts_toggle and not vc_AI:
                                # Check if the bot is in a voice channel
                                if message.guild.voice_client is not None:  
                                    if sync_voice_with_text:
                                        if response_text:
                                            asyncio.create_task(play_tts(message, response_text))
                                        asyncio.create_task(send_message(message.channel, response_text))
                                        return
                                    else:
                                        asyncio.create_task(send_message(message.channel, response_text))
                                        asyncio.create_task(play_tts(message, response_text))
                                        return
                                else:
                                    asyncio.create_task(send_message(message.channel, response_text))
                                    return
                            else:
                                asyncio.create_task(send_message(message.channel, response_text))
                                return

                    break  # Exit the loop if everything is successful
                                
                except Exception as e:
                    error_message = str(e)
                    add_to_history("Error", error_message)

                    if "500" in error_message:
                        print(e)
                    elif "503" in error_message:
                        add_to_history("Error", "The service may be temporarily overloaded or down. Please try again later.")
                        print(f"Error generating content: {e}")
                        await message.channel.send("The service may be temporarily overloaded or down. Please try again later.")
                        break
                    elif "403" in error_message:
                        add_to_history("Error", "Your API key doesn't have the required permissions.")
                        await message.channel.send("Your API key has denied permissions.")
                        print(f"Error generating content: {e}")
                        break
                    elif "504" in error_message:
                        add_to_history("Error", "The service is unable to finish processing within the deadline.")
                        await message.channel.send("The service is unable to finish processing within the deadline.")
                        print(f"Error generating content: {e}")
                        break
                    elif "429" in error_message:
                        add_to_history("Error", "The service is being rate limited.")
                        await message.channel.send("You've exceeded the rate limit, Please try again later.")
                        print(f"Error generating content: {e}")
                        break
                    else:
                        print(f"Error generating content: {e}")
                        add_to_history("Error", f"Error generating content: {str(e)}")
                        await message.channel.send("An error occurred while generating the response.")
                        break
                except TimeoutError:
                    print(f"Error: Server timeout.")
                    add_to_history("Error", "Server timeout.")
                    await message.channel.send("Error: Server timeout.")
                    break
                except discord.errors.HTTPException as e:
                    if e.status == 429:  # Rate limited
                        await message.channel.send("Whoa there, slow down! I'm being rate limited.")
                        # Optionally: log the error or implement retry logic
                    else:
                        raise  # Re-raise other exceptions
                    break
    else:
        await bot.process_commands(message)




# Add other commands and event handlers as needed

ai_toggle_per_channel = {}


# Function to find a voice channel with the specified name, ignoring emojis
def find_voice_channel(guild, channel_name):
    # Remove emojis from the provided channel name
    stripped_name = re.sub(r'[^\x00-\x7F]+', '', channel_name).strip()
    
    # Check if there is an exact match
    voice_channel = get(guild.voice_channels, name=stripped_name)
    
    if voice_channel:
        return voice_channel
    
    # If no exact match, check against channels that may match the stripped name
    for channel in guild.voice_channels:
        if re.sub(r'[^\x00-\x7F]+', '', channel.name) == stripped_name:
            return channel
    
    return None

import traceback
from discord.ext import voice_recv
from io import BytesIO
from pydub import AudioSegment

class VoiceListener(commands.Cog):
    from discord import app_commands
    global ffmpeg_executable_path
    SILENCE_THRESHOLD = -40  # Threshold in dB
    SILENCE_DURATION = 2  # Duration in seconds to consider as silence
    ffmpeg_path = ffmpeg_executable_path

    def __init__(self, bot):
        self.bot = bot
        self.recording = False
        self.audio_data = BytesIO()
        self.last_sound_time = time.time()
        self.silence_timeout = 2.5  # 2.5 seconds of silence
        self.current_user = None  # Initialize current user
        self.channel = None  # Initialize channel attribute
        self.message = None  # Initialize message attribute
        self.check_silence_task = asyncio.create_task(self.check_silence())
        
    
    
    @commands.command()
    async def join(self, ctx):
        if ctx.author.voice:
            vc = await ctx.author.voice.channel.connect(cls=voice_recv.VoiceRecvClient)
            vc.listen(voice_recv.BasicSink(self.voice_callback))
            self.channel = ctx.channel  # Store the channel reference here
            self.message = ctx.message  # Store the message object
            await ctx.send("Joined the voice channel.")
        else:
            await ctx.send("You need to be in a voice channel to use this command.")

    def voice_callback(self, user, data: voice_recv.VoiceData):
        try:
            current_time = time.time()
            audio_segment = AudioSegment(data=data.pcm, sample_width=2, frame_rate=48000, channels=2)
            volume = audio_segment.dBFS

            if volume > -35:  # Adjust as needed
                if not self.recording:
                    self.start_recording(user)
                self.last_sound_time = current_time
                self.audio_data.write(data.pcm)
        except Exception as e:
            print(f"Error in voice callback: {e}")
            traceback.print_exc()

    async def check_silence(self):
        while True:
            await asyncio.sleep(0.1)  # Check every 100ms
            current_time = time.time()

            if self.recording and (current_time - self.last_sound_time > self.silence_timeout):
                await self.stop_recording()

    def start_recording(self, user):
        self.current_user = user  # Set the current user for recording
        self.recording = True
        self.audio_data = io.BytesIO()
        print(f"Started recording for {self.current_user}...")

    async def stop_recording(self):
        try:
            if self.recording and self.current_user:
                self.recording = False
                print(f"Stopped recording for {self.current_user} due to silence.")

                if self.audio_data.tell() > 0:
                    audio_segment = AudioSegment(
                        data=self.audio_data.getvalue(),
                        sample_width=2,
                        frame_rate=48000,
                        channels=2
                    )

                    silence_duration = 1000
                    audio_segment = AudioSegment.silent(duration=silence_duration).append(audio_segment, crossfade=0)
                    user_filename = f"{self.current_user}_prompt.mp3"

                    try:
                        audio_segment.export(user_filename, format="mp3")
                        print(f"Recording saved as {user_filename}")

                        gemini_file = genai.upload_file(path=user_filename)
                        response = chat_session.send_message([f"{self.current_user}'s VoiceChat prompt(DONT REPLY TO THIS INSTRUCTION): ", gemini_file])
                        response_text = response.text.strip()

                        # Pass interaction and user object to play_tts
                        async def generate_tts(message, response_text):
                                global vc_voice
                                ffmpeg_path = ffmpeg_executable_path
                                VOICES = [
                                    'en-US-BrianNeural',     # 1
                                    'en-US-JennyNeural',     # 2
                                    'en-US-GuyNeural',       # 3
                                    'en-GB-SoniaNeural',     # 4
                                    'en-AU-NatashaNeural',   # 5
                                    'en-IN-NeerjaNeural',    # 6
                                    'en-NZ-MitchellNeural',  # 7
                                    'en-CA-ClaraNeural',     # 8
                                    'en-IE-EmilyNeural',     # 9
                                    'en-SG-WayneNeural',     # 10
                                    'en-ZA-LeonNeural',      # 11
                                    'en-GB-RyanNeural',      # 12
                                    'en-AU-WilliamNeural',   # 13
                                    'en-IN-PrabhatNeural',   # 14
                                    'en-NZ-MollyNeural',     # 15
                                    'en-CA-LiamNeural',      # 16
                                    'en-IE-OrlaNeural',      # 17
                                    'en-SG-LunaNeural',      # 18
                                    'en-US-AriaNeural',      # 19
                                    'en-GB-MaisieNeural'     # 20
                                ]
                                print(f"Prompt before filtering: ({response_text})")

                                if os.path.exists(ffmpeg_path):
                                    TEXT = response_text
                                    TEXT = re.sub(r'<.*?>', '', TEXT)  #Remove < > tags only

                                    #Handle Emoji Only. (Instead of deleting other text)
                                    TEXT = re.sub(r'[^\w\s\.,!?"\':;-]', '', TEXT)

                                    TEXT = TEXT.strip()

                                    if isinstance(TEXT, bytes):
                                        TEXT = TEXT.decode('utf-8', errors='ignore')

                                    if not isinstance(TEXT, str):
                                        print(f"generate_tts ERROR: Text must be a string, not a {type(TEXT)} object. Check text contents before conversion")
                                        return False
                                    print(f"Prompt after filtering: ({TEXT})")


                                    voice_index = vc_voice - 1
                                    if voice_index < 0 or voice_index >= len(VOICES):
                                        voice_index = 0

                                    VOICE = VOICES[voice_index]
                                    print(f"Chosen voice: {VOICE}, type: {type(VOICE)}")
                                    OUTPUT_FILE = "system/RAM/vc/Generated_VC.wav"

                                    try:
                                        print(f"Generating TTS for: {TEXT}")
                                        communicate = edge_tts.Communicate(TEXT, VOICE)
                                        print(f"communicate object: {communicate}")
                                        print("Saving TTS...")
                                        await communicate.save(OUTPUT_FILE)
                                        print("TTS file saved!")

                                        file_size = os.path.getsize(OUTPUT_FILE)
                                        print(f"Generated file size: {file_size} bytes")
                                        if file_size == 0:
                                            print("Warning: Generated file is empty.")
                                        return True

                                    except Exception as e:
                                        print(f"Error during TTS generation: {e}")
                                        return False

                                else:
                                    print("Failed to generate VC. FFmpeg not found.")
                                    await message.channel.send("Failed to generate VC. FFmpeg not found.")

                        async def play_tts(message, response_text):
                            ffmpeg_path = ffmpeg_executable_path

                            if os.path.exists(ffmpeg_path):
                                if not message.guild.voice_client:
                                    print("You need to be in a voice channel to play the TTS.")
                                else:
                                    await generate_tts(message, response_text)
                                    OUTPUT_FILE = "system/RAM/vc/Generated_VC.wav"

                                    # Check if the author is in a voice channel
                                    if message.user.voice:
                                        voice_channel = message.user.voice.channel
                                        voice_client = message.guild.voice_client

                                        # Connect to the voice channel if not already connected
                                        if not voice_client:
                                            voice_client = await voice_channel.connect()

                                        # Play the audio file if no audio is playing
                                        if not voice_client.is_playing():
                                            voice_client.play(discord.FFmpegPCMAudio(OUTPUT_FILE, executable=ffmpeg_path))
                                        else:
                                            voice_client.play(discord.FFmpegPCMAudio(OUTPUT_FILE, executable=ffmpeg_path))
                            else:
                                print("Failed to generate VC. FFmpeg not found.")
                                await message.channel.send("Failed to generate VC. FFmpeg not found.")

                        await play_tts(self.message, response_text)
                        print(f"Response sent: {response_text}")

                    except Exception as e:
                        print(f"Error exporting audio or sending message: {e}")
                        traceback.print_exc()
                else:
                    print("No audio data found.")
            else:
                print("Recording already stopped or no user set.")
        except Exception as e:
            print(f"Error stopping recording: {e}")
            traceback.print_exc()

    @commands.command()
    async def leave(self, ctx):
        if ctx.voice_client:
            await ctx.voice_client.disconnect()
            await ctx.send("Left the voice channel.")
    
    @app_commands.command(name="vc", description="Voice Chat Commands")
    @app_commands.describe(action="Action to perform", channel_name="Voice channel (optional)", voice_number="Voice number (1-20)")
    @app_commands.choices(action=[
        app_commands.Choice(name="join", value="join"),
        app_commands.Choice(name="leave", value="leave"),
        app_commands.Choice(name="status", value="status"),
        app_commands.Choice(name="voice", value="voice"),
        app_commands.Choice(name="replay", value="replay")
    ])
    async def vc(self, interaction: discord.Interaction, action: str, channel_name: str = None, voice_number: int = None):
        await interaction.response.defer()
        member = interaction.user

        try:
            if action == "join":
                if not channel_name:  # No channel name provided
                    if member.voice:  # Check if user is in a voice channel
                        voice_channel = member.voice.channel
                    else:
                        await interaction.followup.send("You are not in a voice channel.", ephemeral=True)
                        return
                else:  # Channel name provided
                    voice_channel = discord.utils.get(interaction.guild.voice_channels, name=channel_name)
                    if not voice_channel:
                        await interaction.followup.send(f"Voice channel '{channel_name}' not found.", ephemeral=True)
                        return

                # Check if bot is already connected to a voice channel in current guild
                if interaction.guild.voice_client:
                    if interaction.guild.voice_client.channel == voice_channel:  # Already in the same channel
                        await interaction.followup.send("I'm already in that voice channel.", ephemeral=True)
                        return
                    else:  # Move to the specified channel
                        await interaction.guild.voice_client.move_to(voice_channel)  # Moves bot to correct channel.
                        self.voice_client = interaction.guild.voice_client # Set self.voice_client
                        await interaction.followup.send(f"Moved to {voice_channel.name}.")
                        return

                try:
                    if vc_AI:
                        self.voice_client.listen(voice_recv.BasicSink(self.voice_callback))
                    self.voice_client = await voice_channel.connect(cls=voice_recv.VoiceRecvClient)
                    self.channel = interaction.channel
                    self.message = interaction  # Store interaction
                    await interaction.followup.send(f"Joined {voice_channel.name}.")
                except Exception as e:
                    print(f"Error joining voice channel: {e}") # Print error for debugging
                    await interaction.followup.send(f"Could not join {voice_channel.name} due to an error.", ephemeral=True)

            elif action == "leave":
                if self.voice_client and self.voice_client.is_connected():
                    await self.voice_client.disconnect()
                    self.voice_client = None # Reset voice client
                    await interaction.followup.send("Left the voice channel.")
                else:
                    await interaction.followup.send("Not connected to a voice channel.", ephemeral=True)
            elif action == "status":
                if self.voice_client and self.voice_client.is_connected():
                    await interaction.followup.send("Connected to a voice channel.")
                else:
                    await interaction.followup.send("Not connected to a voice channel.", ephemeral=True)
            elif action == "voice":
                global vc_voice
                voice_number = int(interaction.data.get("options", [{}])[1].get("value", 0))
                member_name = interaction.user.display_name
                add_to_history(member_name, f"/vc voice {voice_number}")
                voice_access = True
                if voice_access:
                    try:
                        # Language handling
                        default_lang = "en"  # Set your default language here (you can change this dynamically)

                        if default_lang == "en":
                            if voice_number < 1 or voice_number > 20:
                                await interaction.followup.send("Voice channel number must be between 1 and 20.")
                                add_to_history("System", "Voice channel number must be between 1 and 20.")
                            else:
                                try:
                                    vc_voice = voice_number
                                    print(f"Chosen voice number: {voice_number}")
                                    add_to_history("System", f"Chosen voice number: {voice_number}")
                                    await interaction.followup.send("Successfully changed voice.")
                                except Exception as e:
                                    print(f"Error changing voice: {e}")
                                    add_to_history("System", f"Error changing voice: {e}")
                                    await interaction.followup.send(f"Error changing voice: {e}")

                        elif default_lang == "eg": # Egyptian Arabic
                            if voice_number < 1 or voice_number > 2:
                                await interaction.followup.send("رقم قناة الصوت يجب أن يكون بين ١ و ٢.") # Voice channel number must be between 1 and 2
                                add_to_history("System", "رقم قناة الصوت يجب أن يكون بين ١ و ٢.") # Voice channel number must be between 1 and 2
                            else:
                                try:
                                    vc_voice = voice_number
                                    print(f"Chosen voice number: {voice_number}")
                                    add_to_history("System", f"اخترت رقم الصوت: {voice_number}") # Chosen voice number
                                    await interaction.followup.send("تم تغيير الصوت بنجاح.") # Successfully changed voice
                                except Exception as e:
                                    print(f"Error changing voice: {e}") 
                                    add_to_history("System", f"حدث خطأ أثناء تغيير الصوت: {e}") # Error changing voice
                                    await interaction.followup.send(f"حدث خطأ أثناء تغيير الصوت: {e}") # Error changing voice

                        elif default_lang == "ru":  # Russian language support
                            if voice_number < 1 or voice_number > 3:
                                await interaction.followup.send("Номер голосового канала должен быть от 1 до 3.")  # Voice channel number must be between 1 and 3
                                add_to_history("System", "Номер голосового канала должен быть от 1 до 3.")
                            else:
                                try:
                                    vc_voice = voice_number
                                    print(f"Chosen voice number: {voice_number}")
                                    add_to_history("System", f"Выбран номер голоса: {voice_number}")  # Chosen voice number: {voice_number}
                                    await interaction.followup.send("Голос успешно изменен.")  # Successfully changed voice
                                except Exception as e:
                                    print(f"Error changing voice: {e}")
                                    add_to_history("System", f"Ошибка при изменении голоса: {e}")  # Error changing voice
                                    await interaction.followup.send(f"Ошибка при изменении голоса: {e}")  # Error changing voice

                    except ValueError:
                        if default_lang == "en":
                            add_to_history("System", "Voice channel number must be a number.")
                            await interaction.followup.send("Voice channel number must be a number.")
                        elif default_lang == "eg": # Egyptian Arabic
                            add_to_history("System", "رقم قناة الصوت لازم يكون رقم.") # Voice channel number must be a number
                            await interaction.followup.send("رقم قناة الصوت لازم يكون رقم.") # Voice channel number must be a number
                        elif default_lang == "ru":  # Russian
                            add_to_history("System", "Номер голосового канала должен быть числом.")  # Voice channel number must be a number
                            await interaction.followup.send("Номер голосового канала должен быть числом.")  # Voice channel number must be a number
                else:
                    print("Failed to generate VC. FFmpeg not found.")
                    await interaction.followup.send("Failed to change voice. VC Not available")
                    add_to_history("System", "Failed to generate VC. FFmpeg not found.")

            elif action == 'replay':
                voice_client = interaction.guild.voice_client
                if voice_client:
                    audio_file = "system/RAM/vc/Generated_voice.mp3"
                    audio_file_wav = "system/RAM/vc/Generated_voice.wav"
                    if os.path.isfile(audio_file):
                        try:
                            voice_client.play(discord.FFmpegPCMAudio(audio_file, executable=ffmpeg_path))
                            add_to_history(member, '/vc replay')
                            add_to_history("System", 'Replaying Audio...')
                            await interaction.followup.send(f'Replaying Audio...')
                        except Exception as e:
                            await interaction.followup.send(f'Error replaying {audio_file}: {e}')
                    elif os.path.isfile(audio_file_wav):
                        try:
                            voice_client.play(discord.FFmpegPCMAudio(audio_file_wav, executable=ffmpeg_path))
                            add_to_history(member, '/vc replay')
                            add_to_history("System", 'Replaying Audio...')
                            await interaction.followup.send(f'Replaying Audio...')
                        except Exception as e:
                            await interaction.followup.send(f'Error replaying {audio_file_wav}: {e}')
                    else:
                        add_to_history("System", "No previous audio generated to replay.")
                        await interaction.followup.send("No previous audio generated to replay.")
                else:
                    await interaction.followup.send("Not connected to a voice channel.")
            else:
                await interaction.followup.send("Sorry, the action you entered is not valid. try `join` or `tts`")
            
        except Exception as e:
            await interaction.followup.send(f"An error occurred: {str(e)}", ephemeral=True)
            print(f"An error occurred: {e}")

# Function to send random images
async def send_random_image(interaction: discord.Interaction, num_images=1):
    # List all files in the images directory
    image_files = [f for f in os.listdir('system/RAM/search-img/') if os.path.isfile(os.path.join('system/RAM/search-img/', f))]

    if not image_files:
        await interaction.followup.send("Sorry, an error has occurred. Please try again.")
        return

    # Filter image files based on extensions (no changes needed here)
    valid_extensions = ['jpg', 'jpeg', 'png', 'webp']
    valid_image_files = [f for f in image_files if f.split('.')[-1].lower() in valid_extensions]

    if not valid_image_files:
        await interaction.followup.send("Sorry, an error has occurred. Please try again.")
        # Empty the images directory
        for f in image_files:
            os.remove(os.path.join('system/RAM/search-img/', f))
        return

    # Select random images
    selected_images = random.sample(valid_image_files, min(num_images, len(valid_image_files)))

    # Send the selected images
    for image_file in selected_images:
        file_path = os.path.join('system/RAM/search-img/', image_file)
        file = discord.File(file_path)
        await interaction.followup.send(file=file)
        os.remove(file_path)  # Delete the image file after sending

@bot.tree.command(name="search_img", description="Search for images.")
async def search_img(interaction: discord.Interaction, query: str, num_images: int = 1, safe_search: bool = True): # Assuming True is the default for safe_search
    await interaction.response.defer()
    try:
        if safe_search:
            safe_search_set = "active"
        else:
            safe_search_set = "off"
        member_name = interaction.user.display_name

        # Search for images using Google Custom Search API
        search_url = "https://www.googleapis.com/customsearch/v1"
        params = {
            'key': SEARCH_API_KEY,
            'cx': CX,
            'q': query,
            'num': num_images,
            'searchType': 'image',
            'safe': safe_search_set
        }

        response = requests.get(search_url, params=params)
        response.raise_for_status()
        results = response.json().get('items', [])

        download_failed = []
        os.makedirs('system/RAM/search-img/', exist_ok=True)

        for result in results:
            try:
                image_url = result['link']
                image_name = image_url.split('/')[-1].split('?')[0]  # Handle URLs with query parameters

                # Generate a random ID using random integers
                random_id = "".join(str(random.randint(0, 9)) for _ in range(16))  # 16-digit random ID

                # Get the original extension, if any
                original_extension = os.path.splitext(image_name)[1]

                # Add .png if no extension or an invalid one
                if not original_extension or original_extension.lower() not in ['.jpg', '.jpeg', '.png', '.webp']:
                    image_name = f"{random_id}.png"
                else:
                    image_name = f"{random_id}{original_extension}"

                response = requests.get(image_url)
                response.raise_for_status()
                with open(os.path.join('system/RAM/search-img/', image_name), 'wb') as f:
                    f.write(response.content)
            except Exception as e:
                download_failed.append({'image_url': image_url, 'error': str(e)})
                await interaction.followup.send(f"Error downloading image from {image_url}. Please try again later.")
        
        # Send the images
        await send_random_image(interaction, num_images)
        add_to_history(member_name, f"/search_img {query}")
        await interaction.followup.send(f"Successfully searched for image `{query}`!")

    except Exception as e:
        await interaction.followup.send("Error occurred. Please try again later.")
        add_to_history("Failed-Search", f"Error during search: {e}")


@bot.tree.command(name="search", description="Perform a web search.")
async def search(interaction: discord.Interaction, query: str, results: int = 5, site: str = None, safe: bool = safe_search):
    await interaction.response.defer()
    if web_search:
        if google_search_api_verified:
            try:
                if results > 20:
                    await interaction.followup.send(f"Error: `{results} is invalid, Maximum number of results is `20`.")
                    return

                search_results = search_google(query, site=site, num_results=results, safe_search=safe)
                if "Error Searching the web." in search_results:
                    print(f"{Fore.RED + Style.BRIGHT}An Error has occured while searching the web:{Style.RESET_ALL} {search_results}")
                    await interaction.followup.send("Error occurred during the search. Try again later.")
                    return

                async with interaction.channel.typing():
                    # Format response (assuming model generates based on prompt)
                    full_prompt = f"Search Results from the web: {search_results}"
                    search_model = genai.GenerativeModel(model_name="gemini-1.5-flash", system_instruction="You are an AI that gets results from the web and summarizes them. and the results from the web are not from the user, its from you")
                    response = search_model.generate_content(full_prompt)
                    await interaction.followup.send(f"- :mag:  **Successfully the Web:**\n{response.text.strip()}")

            except Exception as e:
                await interaction.followup.send("Error occurred during the search. Try again later.")
                print(e)
        else:
            await interaction.followup.send("Error: `Google Search API key` is not valid.")
            print("Error: Google Search API key is not valid.`")
    else:
        await interaction.followup.send("Error: Web Search is disabled, Cannot search the web...")
        print("Error: Web Search is disabled, Cannot search the web...")

@bot.tree.command(name="search_yt", description="Search for YouTube videos.")
async def search_yt(interaction: discord.Interaction, query: str, results: int = 5):
    await interaction.response.defer()
    if web_search:
        if google_search_api_verified:
            try:
                if results > 20:
                    await interaction.followup.send(f"Error: `{results} is invalid, Maximum number of results is `20`.")
                    return

                search_results = search_google(query, site='https://www.youtube.com/', num_results=results)
                if "Error Searching the web." in search_results:
                    print(f"{Fore.RED + Style.BRIGHT}An Error has occured while searching the web:{Style.RESET_ALL} {search_results}")
                    await interaction.followup.send("Error occurred during the search. Try again later.")
                    return

                async with interaction.channel.typing():
                    full_prompt = f"Youtube Videos from the Web: {search_results}"
                    search_model = genai.GenerativeModel(model_name="gemini-1.5-flash", system_instruction="You are an AI that gets results from the web and summarizes them. and the results from the web are not from the user, its from you")
                    response = search_model.generate_content(full_prompt)
                    await interaction.followup.send(f"- :YT: **Successfully Searched Youtube:**\n{response.text.strip()}")

            except Exception as e:
                await interaction.followup.send("Error occurred during search youtube. Try again later.")
                print(e)
        else:
            await interaction.followup.send("Error: `Google Search API key` is not valid.")
            print("Error: Google Search API key is not valid.`")
    else:
        await interaction.followup.send("Error: Web Search is disabled, Cannot search Youtube...")
        print("Error: Web Search is disabled, Cannot search Youtube...")

@bot.tree.command(name="search_save", description="Search and save results.")
async def search_save(interaction: discord.Interaction, query: str, results: int = 5, site: str = None, safe: bool = safe_search):
    global NAME
    await interaction.response.defer()
    if web_search:
        if google_search_api_verified:
            try:
                if results > 20:
                    await interaction.followup.send(f"Error: `{results} is invalid, Maximum number of results is `20`.")
                    return
                search_results = search_google(query, site=site, num_results=results, safe_search=safe)
                if "Error Searching the web." in search_results:
                    print(f"{Fore.RED + Style.BRIGHT}An Error has occured while searching the web:{Style.RESET_ALL} {search_results}")
                    await interaction.followup.send("Error occurred during the search. Try again later.")
                    return

                save_search(query, search_results)
                full_prompt = f"Search Results from the web: {search_results}"
                search_model = genai.GenerativeModel(model_name="gemini-1.5-flash", system_instruction="You are an AI that gets results from the web and summarizes them. and the results from the web are not from the user, its from you")
                response = search_model.generate_content(full_prompt)
                await interaction.followup.send(f"* :mag:  **Successfully Saved the search results to {NAME} for future needs.**\n{response.text.strip()}")

            except Exception as e:
                await interaction.followup.send("Error occurred. Try again later.")
                add_to_history("Failed-Search", f"Error: {e}")
        else:
            await interaction.followup.send("Error: `Google Search API key` is not valid.")
            print("Error: Google Search API key is not valid.`")
    else:
        await interaction.followup.send("Error: Web Search is disabled, Cannot search the web...")
        print("Error: Web Search is disabled, Cannot search the web...")

@bot.tree.command(name="help", description="Get information about available commands")
async def help_command(interaction: discord.Interaction, command_name: str = None):
    await interaction.response.defer()
    global NAME
    try:
        print(f"Help command invoked. command_name: {command_name}")  # Debug log

        if command_name is None:
            print("No command name provided. Sending main help embed.")  # Debug log
            
            # Create the main help embed
            embeds = []

            # First embed (General Commands and Search & Information)
            embed1 = discord.Embed(
                title=f"**{NAME} Command Directory**",
                description=f"**Unlock the Power of {NAME}. Explore the commands below:**",
                color=discord.Color.from_rgb(20, 120, 200)  # Futuristic blue
            )
            embed1.add_field(name=":speech_balloon:  **Conversation & Fun**", value=" ----------------------------------------------------------", inline=False)
            embed1.add_field(name="**/ai**", value=f"Engage in a conversation with {NAME}.", inline=False)
            embed1.add_field(name="**/joke**", value="Get a random joke to brighten your day.", inline=False)
            embed1.add_field(name="**/aitoggle**", value="Enable or disable AI responses for a channel.", inline=False) 
            embed1.add_field(name="**/lang**", value=f"Change the default language for {NAME}.", inline=False)
            embed1.add_field(name="**/report**", value="Report a bug, issue, or user.", inline=False)
            embed1.add_field(name="**/feedback**", value=f"Provide feedback or suggestions for {NAME}.", inline=False)

            embed1.add_field(name=":mag_right:  **Search & Information**", value=" ----------------------------------------------------------", inline=False)
            embed1.add_field(name="**/search**", value="Search the web for information.", inline=False)
            embed1.add_field(name="**/search_yt**", value="Explore videos on YouTube.", inline=False)
            embed1.add_field(name="**/search_save**", value="Save a web search for later use.", inline=False)
            embeds.append(embed1)

            # Second embed (Creative Tools and Memory Management)
            embed2 = discord.Embed(
                color=discord.Color.from_rgb(20, 120, 200)
            )
            embed2.add_field(name=":art:  **Creative Tools**", value=" ----------------------------------------------------------", inline=False)
            embed2.add_field(name="**/img**", value="Generate stunning images from text prompts.", inline=False)
            embed2.add_field(name="**/music**", value="Create unique music based on your description.", inline=False)
            embed2.add_field(name="**/search_img**", value="Search the web for images.", inline=False)

            embed2.add_field(name=":brain:  **Memory Management**", value=" ----------------------------------------------------------", inline=False)
            embed2.add_field(name="**/reset**", value=f"Clears {NAME}'s memory for the current channel.", inline=False)
            embeds.append(embed2)

            # Third embed (Voice Chat and Bot Management)
            embed3 = discord.Embed(
                color=discord.Color.from_rgb(20, 120, 200)
            )
            embed3.add_field(name=":microphone2:  **Voice Chat**", value=" ----------------------------------------------------------", inline=False)
            embed3.add_field(name="**/vc join**", value="Join the specified or your current voice channel.", inline=False)
            embed3.add_field(name="**/vc leave**", value="Leave the current voice channel.", inline=False)
            embed3.add_field(name="**/vc status**", value="Check the current voice channel status.", inline=False)
            embed3.add_field(name="**/tts [text]**", value="Generate Text to Speech from text.", inline=False)
            embed3.add_field(name="**/vc voice [voice_number]**", value="Change the voice for text-to-speech.", inline=False)
            embed3.add_field(name="**/vc replay**", value="Replay the last generated text-to-speech audio.", inline=False)

            embed3.add_field(name="📃  **Additional Commands**", value=" ----------------------------------------------------------", inline=False)
            embed3.add_field(name="**/view [thing]**", value="View more info about something.", inline=False)
            embeds.append(embed3)

            # Send all embeds
            print("Sending help embed response...")  # Debug log
            await interaction.followup.send(embeds=embeds)
            print("Help embeds sent successfully.")  # Debug log

        else:
            # Defer response if a command_name is provided
            print(f"Command name provided: {command_name}. Deferring response.")  # Debug log
            await interaction.response.defer()

            # Find the command using the command tree
            command = bot.tree.get_command(command_name)
            print(f"Command lookup result: {command}")  # Debug log

            if command is None:
                print(f"Command '{command_name}' not found. Sending error message.")  # Debug log
                await interaction.followup.send(f"Command '{command_name}' not found.")
                return

            # Create embed for specific command help
            embed = discord.Embed(
                title=f"**/{command.name}**",  # Use name for slash commands
                description=f"**Description:** {command.description or 'No description provided.'}",
                color=discord.Color.from_rgb(100, 20, 200)  # A deeper, richer blue
            )

            # Skip adding usage information since slash commands don't have a signature
            embed.set_footer(text="Type /help to return to the command list.")

            # Send the follow-up message after deferring
            print(f"Sending follow-up for command: {command_name}")  # Debug log
            await interaction.followup.send(embed=embed)
            print("Follow-up sent successfully.")  # Debug log

    except Exception as e:
        print(f"An error occurred: {e}")  # Debug log
        await interaction.followup.send("An error occurred while processing your request.")

@bot.tree.command(name="ai", description=f'Chat with {NAME}')
async def ai(interaction: discord.Interaction, *, prompt: str, attachment: discord.Attachment = None):
    await interaction.response.defer()

    # Helper function to handle restricted access notification
    async def send_access_notification(message):
        await interaction.followup.send(
            f"Sorry, for full access to {message}, use `/aitoggle`.",
            ephemeral=True
        )

    while True:
        try:
            member_name = interaction.user.display_name
            add_to_history(member_name, prompt)
            full_prompt = f"{member_name}: {prompt}"
            
            # Generate response from the model
            if attachment and attachment.filename.lower().endswith(('png', 'jpg', 'jpeg')):
                save_path = f"system//RAM//read-img//{attachment.filename}"  
                os.makedirs(os.path.dirname(save_path), exist_ok=True)
                await attachment.save(save_path)
                gemini_file = genai.upload_file(save_path)
                response = model.generate_content([gemini_file, full_prompt, "SYSTEM INSTRUCTION(NOT A USER PROMPT): KEEP YOUR RESPONSE UNDER 2000 CHARACTERS!!! and don’t reply to this instruction"])
                response_text = response.text.strip()

            elif attachment:
                await send_access_notification("full file analysis")
                return
            else:
                response = model.generate_content([full_prompt, "SYSTEM INSTRUCTION(NOT A USER PROMPT): KEEP YOUR RESPONSE UNDER 2000 CHARACTERS!!! and don’t reply to this instruction"])
                response_text = response.text.strip()
            
            # Handle different command tags in the response
            if "/img" in response_text:
                await send_access_notification("image generation")
                return
            elif "/music" in response_text:
                await send_access_notification("music generation")
                return
            elif "/memory_save" in response_text:
                await send_access_notification("saving core memory")
                return
            elif "/search" in response_text:
                await send_access_notification("searching the web")
                return
            elif "/search*yt" in response_text:
                await send_access_notification("YouTube search")
                return
            
            # Send response text or handle length errors
            if len(response_text) > 2000:
                print("Error generating content: Response was too long.")
                add_to_history("Error", "Sorry, to send messages over 2000 characters, use </aitoggle:1294672277278818378>.")
                await send_access_notification("sending messages over 2000 characters")
            else:
                add_to_history_bot("", response_text)
                await interaction.followup.send(response_text)
                if os.path.exists('system/RAM/read-img'):
                    # Iterate over each file in the directory and delete it
                    for filename in os.listdir('system/RAM/read-img'):
                        file_path = os.path.join('system/RAM/read-img', filename)
                        try:
                            if os.path.isfile(file_path) or os.path.islink(file_path):
                                os.unlink(file_path)  # Remove file or symbolic link
                            elif os.path.isdir(file_path):
                                shutil.rmtree(file_path)  # Remove directory and its contents
                        except Exception as e:
                            print(f"Error deleting {file_path}: {e}")
                else:
                    print("Directory 'system/RAM/read-img' does not exist.")
            break
        except Exception as e:
            unnecessary_error = await unnecessary_error(e)
            await debug_error(e, "Generating content", interaction.followup)
            if not unnecessary_error:
                break

@bot.tree.command(name="say", description=f"Make {NAME} say something.")
async def say(interaction: discord.Interaction, say: str, channel_name: str = ""):
    await interaction.response.defer()
    member_name = interaction.user.display_name  # Get the user's display name
    add_to_history(member_name, f"/say {say}")
    target_channel = None

    if channel_name:
        if channel_name.isdigit():  # Check if channel_name is a channel ID
            target_channel = interaction.guild.get_channel(int(channel_name))
        else:
            target_channel = discord.utils.get(interaction.guild.channels, name=channel_name)
        
        if target_channel:
            echoed_message = f"{say}"
            await target_channel.send(echoed_message)
            add_to_history("System", say)
            await interaction.followup.send(f"Message sent to {target_channel.name}.")
        else:
            await interaction.followup.send(f"Channel '{channel_name}' not found.")
            add_to_history("System", f"Channel '{channel_name}' not found.")
    else:
        echoed_message = f"{say}"
        await interaction.followup.send(echoed_message)
        add_to_history("System", say)

@bot.tree.command(name="reset", description="Reset the conversation history for this channel.")
async def reset(interaction: discord.Interaction):
    await interaction.response.defer()
    global conversation_history, chat_session

    # Get the channel's name (including handling of emojis)
    channel_name = interaction.channel.name

    # Check if the channel exists in the conversation history
    if channel_name in conversation_history:
        # Reset the history for this channel
        conversation_history[channel_name] = []
        save_history()  # Save the updated history
        chat_session = model.start_chat(history=[])
        await interaction.followup.send(f"Memory has been reset for {interaction.channel.name}.")
    else:
        await interaction.followup.send(f"No conversation history found for {interaction.channel.name}.")

@bot.tree.command(name="profile", description="Get information about a server member.")
async def profile(interaction: discord.Interaction, member: discord.Member = None):
    member = member or interaction.user  # Default to command user if no member is specified
    embed = discord.Embed(
        title=f"{member}",
        description=f"Here is the info we found for {member}",
        color=discord.Color.blue()
    )
    embed.add_field(name="ID", value=member.id)
    embed.add_field(name="Name", value=member.display_name)
    embed.add_field(name="Created at", value=member.created_at)
    embed.add_field(name="Joined at", value=member.joined_at)
    embed.set_thumbnail(url=member.avatar.url)

    await interaction.response.send_message(embed=embed)
    add_to_history(interaction.user.display_name, f"/profile {member}")
    add_to_history("System", f"Info for {member}: ID {member.id}, Name {member.display_name}, Created at {member.created_at}, Joined at {member.joined_at}")

@bot.tree.command(name="serverinfo", description="Get information about this server.")
async def server_info(interaction: discord.Interaction):
    guild = interaction.guild
    embed = discord.Embed(
        title=f"{guild.name}",
        description=f"Here is the info about {guild.name}",
        color=discord.Color.blue()
    )
    embed.add_field(name="Server ID", value=guild.id)
    embed.add_field(name="Member Count", value=guild.member_count)
    embed.add_field(name="Created at", value=guild.created_at)
    embed.set_thumbnail(url=guild.icon.url if guild.icon else "")

    await interaction.response.send_message(embed=embed)
    add_to_history(interaction.user.display_name, f"/serverinfo")
    add_to_history("System", f"Info about {guild.name}: Server ID {guild.id}, Member Count {guild.member_count}, Created at {guild.created_at}")

@bot.tree.command(name="joke", description="Get a random joke. :D")
async def joke(interaction: discord.Interaction):
    response = httpx.get("https://official-joke-api.appspot.com/random_joke")
    if response.status_code == 200:
        joke_data = response.json()
        await interaction.response.send_message(f"{joke_data['setup']} - {joke_data['punchline']}")
        add_to_history(interaction.user.display_name, "/joke")
        add_to_history("System", f"Joke: {joke_data['setup']} - {joke_data['punchline']}")
    else:
        await interaction.response.send_message("Couldn't fetch a joke at the moment. Try again later!")
        add_to_history(interaction.user.display_name, "/joke")
        add_to_history("System", "Couldn't fetch a joke at the moment. Try again later!")



model_name_model = genai.GenerativeModel(
  model_name="gemini-1.5-flash",
  generation_config=gen_config,
  system_instruction="you are a model name generator and the user will give you models and you will have to get the original name from it and nothing else, dont respond with anything else, only the generated name, example: `User: stabilityai/stable-diffusion-xl-base-1.0, You: Stable Diffusion XL Base 1.0`",
  safety_settings=[
        {"category": "HARM_CATEGORY_DANGEROUS", "threshold": "BLOCK_NONE"},
        {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
        {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
        {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
        {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
  ],
)

if not Image_Model == "stabilityai/stable-diffusion-xl-base-1.0":
    response = model_name_model.generate_content(Image_Model)
    Image_Model_Name = response.text.strip()
    print(f"Generated Image Model Name: {Image_Model_Name} | You need to reinvite the bot once to use the new model")
else:
    Image_Model_Name = "Stable Diffusion XL Base 1.0"


from discord import app_commands, ui

view_options = [
    app_commands.Choice(name="What AI Model is currently selected.", value="view-model"),
    app_commands.Choice(name="Google's Gemini 1.5 Flash Model", value="view-1.5-flash"),
    app_commands.Choice(name="Google's Gemini 1.5 Flash 8B Model", value="view-1.5-flash-8b"),
    app_commands.Choice(name="Google's Gemini 1.5 Pro Model", value="view-1.5-pro"),
    app_commands.Choice(name="Google's Gemini Experimental 1114 Model", value="view-exp-1114"),
    app_commands.Choice(name="Google's Gemini Experimental 1121 Model", value="view-exp-1121"),
    app_commands.Choice(name="Google's Gemini Experimental 1206 Model", value="view-exp-1206"),
    app_commands.Choice(name="Google's LearnLM 1.5 Pro Experimental Model", value="view-learnlm-1.5-pro-exp"),
]

@bot.tree.command(name="view", description='View more info about something / View current setting.')
@app_commands.describe(view="View about...")
@app_commands.choices(view=view_options)
async def view_command(interaction: discord.Interaction, view: str):
    await interaction.response.defer()
    global NAME

    if view == "view-model":
        user = interaction.user.name
        user_settings = get_user_settings(user) # Call the function directly
        selected_model_name = user_settings['model_name']
        await interaction.followup.send(f"Current AI Model: {selected_model_name}")
    elif view == "view-1.5-flash":
        embed = discord.Embed(
            title="Google Gemini 1.5 Flash",
            description=(
                "Gemini 1.5 Flash is a cutting-edge AI model from Google DeepMind, engineered for unparalleled speed, efficiency, and scalability. "
                "It serves as a lightweight alternative to Gemini 1.5 Pro while maintaining advanced capabilities for multimodal reasoning and long-context processing."
            ),
            color=discord.Color.blue()
        )

        embed.set_image(url="https://github.com/user-attachments/assets/f246b746-fddf-4809-a036-9c20f31c67f9")

        embed.add_field(
            name="Key Features",
            value=(
                "**1. Speed and Efficiency:** Tailored for low latency and high throughput, suitable for high-frequency tasks.\n"
                "**2. Contextual Power:** Processes up to 1 million tokens, enabling complex, long-form analysis.\n"
                "**3. Multimodal Expertise:** Handles tasks across text, images, video, and audio with remarkable accuracy.\n"
                "**4. Cost Efficiency:** With an economical cost of $0.0375 per million input tokens, it's ideal for enterprise use.\n"
                "**5. Advanced Training:** Utilizes a distillation process to inherit essential knowledge from Gemini 1.5 Pro."
            ),
            inline=False
        )

        embed.add_field(
            name="Ideal Use Cases",
            value=(
                "- High-volume summarization\n"
                "- Advanced chat applications\n"
                "- Image and video captioning\n"
                "- Extracting data from long documents\n"
                "- Efficient code analysis"
            ),
            inline=False
        )

        embed.add_field(
            name="Recent Updates",
            value=(
                "Gemini 1.5 Flash is available in over 230 countries and supports more than 40 languages. "
                "It has expanded accessibility for developers via Google AI Studio and Vertex AI. "
                "The model is particularly optimized for enterprises requiring rapid processing of large datasets."
            ),
            inline=False
        )

        embed.add_field(
            name="Technical Highlights",
            value=(
                "- **Input Token Limit:** 1,048,576 tokens\n"
                "- **Supported Inputs:** Audio, video, images, and text\n"
                "- **Training Innovation:** Leveraged distillation from larger models to maximize speed and precision"
            ),
            inline=False
        )
        embed.set_footer(
            text=f"Learn more at Google AI Studio and Google AI Documentation.\n\nhttps://aistudio.google.com/\nhttps://deepmind.google/technologies/gemini/flash/\n\nSupported Model for {NAME}"
        )
        await interaction.followup.send(embed=embed)
    elif view == "view-1.5-flash-8b":
        embed = discord.Embed(
            title="Google Gemini 1.5 Flash 8B",
            description=(
                "The Gemini 1.5 Flash 8B is Google's most affordable and compact AI model, "
                "designed for developers seeking high performance at a low cost. It is optimized for "
                "multimodal tasks and long-context processing while maintaining competitive accuracy and speed."
            ),
            color=discord.Color.green()
        )
        embed.set_image(url="https://github.com/user-attachments/assets/efc4ad34-86f1-4cbd-8d31-ae0b71a71346")

        embed.add_field(
            name="Key Features",
            value=(
                "**1. Affordability:** $0.0375 per million input tokens; $0.15 per million output tokens.\n"
                "**2. Performance:** 2x higher rate limits and low latency for small prompts.\n"
                "**3. Multimodal Support:** Handles text, image, audio, and video tasks.\n"
                "**4. Scalability:** Tailored for high-volume tasks like transcription and chatbots.\n"
                "**5. Accessibility:** Available via Google AI Studio and Gemini API."
            ),
            inline=False
        )

        embed.add_field(
            name="Ideal Applications",
            value=(
                "- Real-time chatbots\n"
                "- Language translation\n"
                "- High-volume data analysis\n"
                "- Multimodal input tasks\n"
                "- Long-context summarization"
            ),
            inline=False
        )

        embed.add_field(
            name="Why Choose Gemini 1.5 Flash 8B?",
            value=(
                "This model offers a balance of cost efficiency and functionality, making it ideal "
                "for developers working on scalable AI applications."
            ),
            inline=False
        )
        embed.set_footer(
            text=f"Learn more at Google AI Studio and Google AI Documentation.\n\nhttps://aistudio.google.com/\nhttps://deepmind.google/technologies/gemini/flash/\n\nSupported Model for {NAME}"
        )
        await interaction.followup.send(embed=embed)
    elif view == "view-1.5-pro":
        embed = discord.Embed(
            title="Google Gemini 1.5 Pro",
            description=(
                "Gemini 1.5 Pro is Google's flagship AI model in the Gemini series, "
                "built for advanced performance across multimodal tasks, long-context handling, and efficient cost optimization."
            ),
            color=discord.Color.blue()
        )

        # Add an image
        embed.set_image(url="https://github.com/user-attachments/assets/73f4e79f-a324-440d-a478-d16f9aebfaed")

        # Key Features
        embed.add_field(
            name="Key Features",
            value=(
                "**1. Long Context Window:** Supports up to 2 million tokens, ideal for analyzing large documents, videos, or repositories.\n"
                "**2. Multimodal Capabilities:** Processes and generates text, images, audio, and video seamlessly.\n"
                "**3. Cost Efficiency:** 64% reduction in token input costs and 52% reduction in output costs.\n"
                "**4. Speed & Latency:** Delivers 2x faster output and 3x lower latency compared to earlier versions.\n"
                "**5. Enhanced Reasoning:** Excels in benchmarks like MMLU-Pro and HiddenMath, with significant performance gains."
            ),
            inline=False
        )

        # Ideal Use Cases
        embed.add_field(
            name="Ideal Use Cases",
            value=(
                "- **Document Analysis:** Handles PDFs exceeding 1,000 pages.\n"
                "- **Coding Assistance:** Generates Python code and assists in debugging.\n"
                "- **Creative Applications:** Supports image generation and storytelling.\n"
                "- **Education:** Creates interactive learning tools and personalized teaching aids.\n"
                "- **Customer Support:** Efficient multilingual support for global users."
            ),
            inline=False
        )

        # Benchmark Achievements
        embed.add_field(
            name="Benchmark Achievements",
            value=(
                "**- MMLU-Pro:** Improved scores by ~7%, achieving >85% accuracy across 57 subjects.\n"
                "**- HiddenMath:** ~20% better performance on complex mathematical problems.\n"
                "**- HellaSwag:** Achieved 93.3% accuracy for sentence completion tasks.\n"
                "**- HumanEval:** Scored 84.1% on problem-solving and code generation tasks."
            ),
            inline=False
        )

        # Updates and Enhancements
        embed.add_field(
            name="Recent Enhancements",
            value=(
                "1. New caching strategies to optimize token usage.\n"
                "2. More concise default responses for cost-efficient usage.\n"
                "3. Safer, more user-aligned outputs for diverse applications."
            ),
            inline=False
        )

        # Pricing Information
        embed.add_field(
            name="Pricing Details",
            value=(
                "Gemini 1.5 Pro offers **64% lower input token costs**, making it more accessible for high-volume tasks. "
                "Available through Google AI Studio and the Gemini API."
            ),
            inline=False
        )
        embed.set_footer(
            text=f"Learn more at Google AI Studio and Google AI Documentation.\n\nhttps://aistudio.google.com/\nhttps://deepmind.google/technologies/gemini/pro/\n\nSupported Model for {NAME}"
        )
        await interaction.followup.send(embed=embed)
    elif view == "view-exp-1114":
        embed = discord.Embed(
            title="Google Gemini Experimental 1114",
            description=(
                "Gemini 1114 remains a strong performer in the Gemini series, excelling in creative and multimodal tasks. "
                "However, it faces growing competition from newer models such as Gemini 1121 and GPT-4o, with some limitations in logical reasoning."
            ),
            color=discord.Color.orange()
        )

        # Add an image
        embed.set_image(url="https://github.com/user-attachments/assets/2399fdfe-13f4-4b70-9731-36d6166a74fa")  # Replace with an actual image URL

        # Key Features
        embed.add_field(
            name="Key Features",
            value=(
                "**1. Benchmark Ranking:** Currently in 3rd place in performance, behind GPT-4o and Gemini 1121.\n"
                "**2. Multimodal Strengths:** Excellent at handling tasks involving text and image processing.\n"
                "**3. Slower Response Time:** While capable, it has a slower response time in tasks requiring extensive reasoning.\n"
                "**4. Creative Writing:** Performs well in content generation for creative use cases.\n"
                "**5. Logical Reasoning:** Struggles with some logical tasks, leading to occasional inaccuracies."
            ),
            inline=False
        )

        # Ideal Applications
        embed.add_field(
            name="Ideal Applications",
            value=(
                "- Creative content generation.\n"
                "- Image captioning and multimodal interaction.\n"
                "- Problem-solving tasks requiring synthesis of knowledge.\n"
                "- Text-based applications needing creativity and coherence."
            ),
            inline=False
        )

        # Limitations and Challenges
        embed.add_field(
            name="Limitations and Challenges",
            value=(
                "Despite its strengths, **Gemini 1114** is slower in some contexts and can struggle with logical reasoning, "
                "impacting accuracy in certain tasks."
            ),
            inline=False
        )

        # Current Status
        embed.add_field(
            name="Current Status",
            value=(
                "While **Gemini 1114** remains a strong model, it has been outperformed in certain areas by newer models like **Gemini 1121** and **GPT-4o**."
            ),
            inline=False
        )
        embed.set_footer(
            text=f"Learn more at Google AI Studio and Google AI Documentation.\n\nhttps://aistudio.google.com/\nhttps://deepmind.google/technologies/gemini/\n\nSupported Model for {NAME}"
        )
        await interaction.followup.send(embed=embed)
    elif view == "view-exp-1121":
        embed = discord.Embed(
            title="Google Gemini Experimental 1121",
            description=(
                "Google Gemini-Exp-1121 is the latest iteration in Google's AI models, breaking records in the competitive landscape. "
                "With a focus on multi-turn dialogue, reasoning, and enhanced visual understanding, Gemini-Exp-1121 outperforms GPT-4o in recent benchmarks, "
                "setting a new standard for coding, problem-solving, and multimodal AI applications."
            ),
            color=discord.Color.purple()
        )

        # Set image for the embed (use the actual URL to a relevant image or logo)
        embed.set_image(url="https://github.com/user-attachments/assets/2399fdfe-13f4-4b70-9731-36d6166a74fa")  # Replace with an actual image URL

        # Key Features
        embed.add_field(
            name="Key Features",
            value=(
                "**1. Top Ranking in AI Benchmarks:** Surpassed GPT-4o to lead the Chatbot Arena leaderboard.\n"
                "**2. Enhanced Coding Performance:** Optimized for more complex programming tasks with high accuracy.\n"
                "**3. Stronger Reasoning Abilities:** Capable of handling multi-step problem-solving with ease.\n"
                "**4. Visual Understanding:** Exceptional at processing visual inputs, including images and video.\n"
                "**5. Multi-Turn Dialogue Excellence:** Can maintain context in long conversations, excelling in complex dialogues."
            ),
            inline=False
        )

        # Ideal Use Cases
        embed.add_field(
            name="Ideal Use Cases",
            value=(
                "- Enterprise applications and development tasks.\n"
                "- Coding assistance and code generation.\n"
                "- Complex problem-solving in academic and research settings.\n"
                "- Multimodal tasks such as image analysis, text generation, and video processing."
            ),
            inline=False
        )

        # Status and Performance
        embed.add_field(
            name="Performance and Current Status",
            value=(
                "Gemini-Exp-1121 has not only surpassed its predecessors but also remains ahead of GPT-4o in benchmarks. "
                "It continues to demonstrate cutting-edge advancements in AI, especially for tasks requiring deep reasoning and creative outputs."
            ),
            inline=False
        )

        # Limitations and Challenges
        embed.add_field(
            name="Limitations",
            value=(
                "The key limitation of **Gemini-Exp-1121** is its 32k token context window. While this is highly capable, it pales in comparison to the "
                "2 million token context window of **Gemini 1.5 Pro** or the 1 million tokens offered by **Gemini 1.5 Flash**. This relatively smaller context "
                "window may limit its ability to handle extremely long documents or multi-turn dialogues requiring very large context retention. "
                "However, this limitation is expected to be addressed and improved in future updates, potentially increasing the context window significantly."
            ),
            inline=False
        )

        embed.set_footer(
            text=f"Learn more at Google AI Studio and Google AI Documentation.\n\nhttps://aistudio.google.com/\nhttps://deepmind.google/technologies/gemini/\n\nSupported Model for {NAME}"
        )
        await interaction.followup.send(embed=embed)
    elif view == "view-learnlm-1.5-pro-exp":
        embed = discord.Embed(
            title="LearnLM 1.5 Pro Experimental",
            description=(
                "LearnLM 1.5 Pro Experimental is a task-specific model developed by Google, aimed at revolutionizing the learning experience. "
                "By adhering to learning science principles, it supports active learning, adapts to student needs, and fosters curiosity."
            ),
            color=discord.Color.light_grey()
        )

        embed.set_image(url="https://github.com/user-attachments/assets/5631a182-654f-48e5-b82f-770bb0ac74ae")

        embed.add_field(
            name="Key Features",
            value=(
                "**1. Active Learning:** Encourages students to engage actively with the material and reflect on their thought process.\n"
                "**2. Adaptivity:** Adjusts the difficulty of tasks based on the student's performance and goals.\n"
                "**3. Cognitive Load Management:** Structures information for easier absorption, using multiple modalities.\n"
                "**4. Stimulating Curiosity:** Fosters a positive learning environment to inspire motivation.\n"
                "**5. Metacognition:** Helps students monitor their progress and make necessary adjustments."
            ),
            inline=False
        )

        embed.add_field(
            name="Use Cases",
            value=(
                "- Test Preparation\n"
                "- Concept Teaching\n"
                "- Simplifying Complex Texts for Different Learning Levels\n"
                "- Helping Students Reflect on Their Learning Journey"
            ),
            inline=False
        )

        embed.add_field(
            name="Limitations",
            value=(
                "While LearnLM 1.5 Pro Experimental excels in personalized learning tasks, it is still evolving and might not be as robust "
                "in more general AI tasks outside of educational contexts. Future updates are expected to enhance its capabilities."
            ),
            inline=False
        )

        embed.set_footer(
            text=f"Learn more at Google AI Studio and Google AI Documentation.\n\nhttps://aistudio.google.com/\nhttps://ai.google.dev/gemini-api/docs/learnlm\n\nSupported Model for {NAME}"
        )

        await interaction.followup.send(embed=embed)
    elif view == "view-exp-1206":
        embed = discord.Embed(
            title="Google Gemini Experimental 1206",
            description=(
                "Google Gemini Experimental 1206 is the latest and most powerful AI model from Google DeepMind, "
                "surpassing all previous LLMs and AI models in terms of performance, capabilities, and versatility. "
                "While it is incredibly advanced across many domains, its **2 million tokens context window** allows it to "
                "handle larger and more complex tasks than ever before. This includes a wide variety of tasks, ranging from coding to "
                "scientific research, and general problem-solving. Its multimodal capabilities — processing text, images, audio, and video "
                "seamlessly — represent a new era in AI performance."
            ),
            color=discord.Color.from_rgb(216, 164, 68)
        )
        embed.set_image(url="https://github.com/user-attachments/assets/2399fdfe-13f4-4b70-9731-36d6166a74fa")  # Replace with actual image URL for visual appeal

        embed.add_field(
            name="Key Features",
            value=(
                "**1. Over 2 million tokens context window:** Allows for handling of massive datasets and long-term reasoning across various domains.\n"
                "**2. Unmatched multimodal capabilities:** Processes text, images, audio, and video with seamless coherence, significantly outperforming previous models.\n"
                "**3. Record-breaking performance across multiple benchmarks:** Achieves exceptional scores in reasoning, math, scientific research, and even code generation.\n"
                "**4. Superior code generation and review abilities:** Outperforms models like GPT-4 in programming, offering complete solutions for web apps, code reviews, and more.\n"
                "**5. Advanced reasoning across complex domains:** Solves multi-layered problems in fields like mathematics, physics, law, and medicine at an unprecedented level of sophistication.\n"
                "**6. Next-gen contextual understanding:** Outperforms predecessors with its ability to handle longer and more complex inputs with accuracy."
            ),
            inline=False
        )

        embed.add_field(
            name="Limitations",
            value=(
                "While **Gemini 1206** represents a significant leap in AI capabilities, it does face some challenges, particularly in the vision space, "
                "where it currently ranks **4th** on the leaderboard, behind **Gemini 1121**. However, this gap is expected to close in future updates, "
                "as its development is ongoing and vision capabilities are a priority for enhancement."
            ),
            inline=False
        )

        embed.add_field(
            name="Performance in Vision Tasks",
            value=(
                "While **Gemini 1206** outperforms all LLMs in natural language processing and reasoning tasks, it ranks **4th** in the vision leaderboard, "
                "behind **Gemini 1121**. Nevertheless, with advancements expected in upcoming versions, this limitation is likely to be addressed, "
                "making **Gemini 1206** an all-around leader in AI capabilities."
            ),
            inline=False
        )

        embed.add_field(
            name="Why Gemini 1206 is Revolutionary",
            value=(
                "Gemini 1206 is a paradigm shift in the world of AI. Its multimodal abilities extend beyond the traditional text-based understanding, "
                "integrating various forms of media into a unified cognitive system. With its **2 million token context window**, the model can manage "
                "longer, more complex tasks, making it ideal for industries such as **finance, healthcare, research**, and **software development**. "
                "Its performance across **coding** (including multiple programming languages) and **scientific research** surpasses current models by a wide margin, "
                "demonstrating its versatility and adaptability in real-world applications."
            ),
            inline=False
        )
        embed.set_footer(
            text=f"Learn more at Google AI Studio and Google AI Documentation.\n\nhttps://aistudio.google.com/\nhttps://deepmind.google/technologies/gemini/\n\nSupported Model for {NAME}"
        )
        await interaction.followup.send(embed=embed)




# Define the list of models as choices (name, value)
image_model_choices = [
    app_commands.Choice(name=f"{Image_Model_Name} (Default)", value=Image_Model),
    app_commands.Choice(name="Stable Diffusion 3 Medium Diffusers", value="stabilityai/stable-diffusion-3-medium-diffusers"),
    app_commands.Choice(name="DALL-E 3 XL V2", value="ehristoforu/dalle-3-xl-v2"),
    app_commands.Choice(name="FLUX.1 Schnell", value="black-forest-labs/FLUX.1-schnell"),
    app_commands.Choice(name="FLUX Anime 2", value="dataautogpt3/FLUX-anime2"),
    app_commands.Choice(name="Chip & DallE", value="Yntec/Chip_n_DallE"),
    app_commands.Choice(name="Flux.1 DEV", value="black-forest-labs/FLUX.1-dev"),
    app_commands.Choice(name="Flux.1 DEV LoRA Art", value="Shakker-Labs/FLUX.1-dev-LoRA-Garbage-Bag-Art"),
    app_commands.Choice(name="Flux.1 DEV LoRA Playful Metropolis Art", value="Shakker-Labs/FLUX.1-dev-LoRA-playful-metropolis"),
    app_commands.Choice(name="Flux.1 DEV LoRA Logo Design (Create logos)", value="Shakker-Labs/FLUX.1-dev-LoRA-Logo-Design"),
    app_commands.Choice(name="Flux.1 DEV LoRA Add Details (Advanced Details)", value="Shakker-Labs/FLUX.1-dev-LoRA-add-details"),
]

@bot.tree.command(name="img", description="Generate an image based on your prompt.")
@app_commands.describe(prompt="The image prompt", model="Choose a model to generate the image (optional)")
@app_commands.choices(model=image_model_choices)
async def img(interaction: discord.Interaction, prompt: str, model: str = None):
    if HUGGING_FACE_API == "YOUR_HUGGING_FACE_API_KEY":
        await interaction.followup.send("Sorry, You have entered an Invalid Hugging Face API Key to use `/img`!") 
        return

    else:
        await interaction.response.defer()  # Defer the response to allow for processing time
        is_nsfw = False

        if prompt:
            check_prompt_response_text = prompt

            # Custom view for mod actions with buttons
            class ConfirmUnbanView(ui.View):
                def __init__(self, user, message):
                    super().__init__(timeout=None)
                    self.user = user
                    self.message = message

                @ui.button(label="Unban", style=discord.ButtonStyle.success)
                async def confirm_unban_button(self, interaction: discord.Interaction, button: ui.Button):
                    try:
                        # Attempt to unban the user
                        await interaction.guild.unban(self.user, reason="Unbanned by moderator.")
                        
                        # Create an embed to confirm the unban action
                        embed_unban = discord.Embed(
                            title="⚠️ Unban Successful",
                            description=f"{self.user.mention} has been successfully unbanned.",
                            color=discord.Color.green(),
                        )
                        embed_unban.add_field(name="User", value=f"{self.user.mention}\n(ID: {self.user.id})")
                        embed_unban.set_thumbnail(url=self.user.avatar.url)
                        embed_unban.set_footer(text="⚙️ Automated Moderation System")

                        # Send the unban message to the specified '✅・unbanned-users' channel
                        banned_users_channel = discord.utils.get(interaction.guild.channels, name="✅・unbanned-users")
                        if banned_users_channel:
                            await banned_users_channel.send(embed=embed_unban)
                            try:
                                await self.user.send(f"You have been unbanned from {interaction.guild.name}! Here is a new invite: [Invite Link](https://discord.gg/Va8kH3X5gz)")
                            except discord.Forbidden:
                                pass  # If the user has DMs disabled, ignore this error.
                        else:
                            print("The '✅・unbanned-users' channel was not found. Please create it or check the channel name.")

                        # Optionally, edit the original message or remove buttons after unban
                        await message.edit(content=f"{self.user.mention} has been unbanned.", view=None)
                        await interaction.followup.send(f"{self.user.mention} has been unbanned.", ephemeral=True)

                    except discord.NotFound:
                        # If the user is not found in the ban list (already unbanned or not banned)
                        await interaction.followup.send(f"{self.user.mention} is not banned.", ephemeral=True)
                    except discord.Forbidden:
                        # If the bot lacks permission to unban
                        await interaction.followup.send("I do not have permission to unban this user.", ephemeral=True)
                    except Exception as e:
                        # Catch-all for any other exceptions
                        await interaction.followup.send(f"An error occurred: {str(e)}", ephemeral=True)
                @ui.button(label="Cancel", style=discord.ButtonStyle.danger)
                async def cancel_button(self, interaction: discord.Interaction, button: ui.Button):
                    # If cancelled, just delete the confirmation message
                    await interaction.followup.send("Unban action cancelled.", ephemeral=True)
                    await message.edit(content="Unban action cancelled.", view=None)
            class BanDurationSelect(ui.Select):
                def __init__(self, user, message):
                    options = [
                        discord.SelectOption(label='1 Hour', description='Ban the user for 1 hour', value='1h'),
                        discord.SelectOption(label='1 Day', description='Ban the user for 1 day', value='1d'),
                        discord.SelectOption(label='1 Week', description='Ban the user for 1 week', value='1w'),
                        discord.SelectOption(label='Permanent', description='Permanently ban the user', value='permanent')
                    ]
                    super().__init__(placeholder='Select a ban duration...', options=options)
                    self.user = user
                    self.message = message

                async def callback(self, interaction: discord.Interaction):
                    # Defer the interaction to allow time for the operation
                    await interaction.response.defer(ephemeral=True)

                    # Ban the user based on the selected duration
                    if self.values[0] == 'permanent':
                        await interaction.guild.ban(self.user, reason='Permanent ban by moderator.')
                        ban_message = f"{self.user.mention} has been permanently banned."
                    else:
                        await interaction.guild.ban(self.user, reason=f"Banned for {self.values[0]} by moderator.")
                        ban_message = f"{self.user.mention} has been banned for {self.values[0]}."

                    # After banning, show the unban button
                    view = ui.View()
                    unban_button = UnbanButton(self.user, self.message)
                    view.add_item(unban_button)

                    # Update the original message to show the ban result and unban button
                    try:
                        await self.message.edit(content=ban_message, view=view)
                    except discord.NotFound:
                        print("Error: The message to edit no longer exists (was deleted).")

                    # Optionally send a final response (since we deferred earlier)
                    await interaction.followup.send(ban_message, ephemeral=True)
            class UnbanButton(ui.Button):
                def __init__(self, user, message):
                    super().__init__(label="Unban", style=discord.ButtonStyle.success)
                    self.user = user
                    self.message = message

                async def callback(self, interaction: discord.Interaction):
                    # Send a confirmation message asking the moderator if they want to unban the user
                    embed_confirm = discord.Embed(
                        title="Confirm Unban",
                        description=f"Do you want to unban {self.user.mention}?",
                        color=discord.Color.orange()
                    )
                    embed_confirm.set_thumbnail(url=self.user.avatar.url)
                    embed_confirm.set_footer(text="⚠️ Automated Moderation System")

                    # Create a confirmation view with the unban button
                    confirm_view = ConfirmUnbanView(self.user, message)
                    
                    # Send the confirmation message with the unban and cancel buttons
                    await interaction.followup.send(embed=embed_confirm, view=confirm_view, ephemeral=True)
            class ModActionView(ui.View):
                def __init__(self, user, mod_channel, user_channel):
                    super().__init__(timeout=None)
                    self.user = user
                    self.mod_channel = mod_channel
                    self.user_channel = user_channel

                @ui.button(label="Ban", style=discord.ButtonStyle.danger)
                async def ban_button(self, interaction: discord.Interaction, button: ui.Button):
                    # Show ban duration select
                    view = ui.View()
                    select = BanDurationSelect(self.user, interaction.message)  # Use interaction.message
                    view.add_item(select)

                    # Send a message with the ban duration select view
                    await interaction.followup.send("Select the ban duration:", view=view, ephemeral=True)

                    # Do not delete the message as we are still using it
                    # Use interaction.message for further reference

                @ui.button(label="Kick", style=discord.ButtonStyle.primary)
                async def kick_button(self, interaction: discord.Interaction, button: ui.Button):
                    # Kick the user
                    await interaction.guild.kick(self.user, reason="Kicked by moderator due to inappropriate prompt.")
                    await interaction.followup.send(f"{self.user.mention} has been kicked.", ephemeral=True)

                    # Log the kick action to the mod channel
                    await self.mod_channel.send(f"{self.user.mention} was kicked by {interaction.user.mention}.")
                    
                    # Update the message to remove buttons since the action was taken
                    try:
                        await interaction.message.edit(content=f"{self.user.mention} has been kicked.", view=None)
                    except discord.NotFound:
                        print("Message not found for editing after kick.")

                @ui.button(label="Ignore", style=discord.ButtonStyle.secondary)
                async def ignore_button(self, interaction: discord.Interaction, button: ui.Button):
                    try:
                        await message.delete()
                    except Exception as error_message_ignore:
                        print(f"Error on ignore button: {error_message_ignore}")

            try:
                try:
                    model_check_prompt_pro = genai.GenerativeModel(model_name="gemini-1.5-pro-latest", generation_config=gen_config, system_instruction="your purpose is to only give out numbers to the prompts: if the prompt is a different language than English only say `1` and if it's inappropriate only say `2` but dont be too strict with it and if it's corrupted only say `3` and only say `4` if it's good to go and ONLY say `5` if its REALLY inappropriate or smth and needs to be sent to the moderators to temp ban or timeout the user.", safety_settings=sys_security)
                    response_check = model_check_prompt_pro.generate_content(check_prompt_response_text)
                    response_check_text_check = response_check.text.strip()
                    print("Used model Pro on check")
                except Exception as e3:
                    try:
                        model_check_prompt_advanced = genai.GenerativeModel(model_name="gemini-1.5-pro-002", generation_config=gen_config, system_instruction="your purpose is to only give out numbers to the prompts: if the prompt is a different language than English only say `1` and if it's inappropriate only say `2` but dont be too strict with it and if it's corrupted only say `3` and only say `4` if it's good to go and ONLY say `5` if its REALLY inappropriate or smth and needs to be sent to the moderators to temp ban or timeout the user.", safety_settings=sys_security)
                        response_check = model_check_prompt_advanced.generate_content(check_prompt_response_text)
                        response_check_text_check = response_check.text.strip()
                        print("Used Model Pro Advanced on check")
                    except Exception as e2:
                        try:
                            model_check_prompt_flash = genai.GenerativeModel(model_name="gemini-1.5-flash-latest", generation_config=gen_config, system_instruction="your purpose is to only give out numbers to the prompts: if the prompt is a different language than English only say `1` and if it's inappropriate only say `2` but dont be too strict with it and if it's corrupted only say `3` and only say `4` if it's good to go and ONLY say `5` if its REALLY inappropriate or smth and needs to be sent to the moderators to temp ban or timeout the user.", safety_settings=sys_security)
                            response_check = model_check_prompt_flash.generate_content(check_prompt_response_text)
                            response_check_text_check = response_check.text.strip()
                            print("Used Model Flash on check")
                        except Exception as e:
                            print(f"Failed to run all Models for prompt check. | ERROR: {e}")
                print(f"Check: {response_check_text_check}")

                # Handle different check results (1 for translation, 2 for inappropriate, etc.)
                if "1" in response_check_text_check:
                    model_translation_gen_prompt_pro = genai.GenerativeModel(model_name="gemini-1.5-pro-latest", generation_config=gen_config, system_instruction="your porpuse is to only to translate any user's prompt language to english, and nothing else, you must not say anything unless its the translated prompt, just like google translate!", safety_settings=sys_security)
                    model_translation_gen_prompt_advanced = genai.GenerativeModel(model_name="gemini-1.5-pro-002", generation_config=gen_config, system_instruction="your porpuse is to only to translate any user's prompt language to english, and nothing else, you must not say anything unless its the translated prompt, just like google translate!", safety_settings=sys_security)
                    model_translation_gen_prompt_flash = genai.GenerativeModel(model_name="gemini-1.5-flash-latest", generation_config=gen_config, system_instruction="your porpuse is to only to translate any user's prompt language to english, and nothing else, you must not say anything unless its the translated prompt, just like google translate!", safety_settings=sys_security)
                    try:
                        response_translation_gen_image = model_translation_gen_prompt_pro.generate_content(check_prompt_response_text)
                        response_translation_gen_image_text_translate = response_translation_gen_image.text.strip()
                        prompt = f"{response_translation_gen_image_text_translate}"
                        print("Used model Pro")
                    except Exception as e3:
                        try:
                            response_translation_gen_image = model_translation_gen_prompt_advanced.generate_content(check_prompt_response_text)
                            response_translation_gen_image_text_translate = response_translation_gen_image.text.strip()
                            prompt = f"{response_translation_gen_image_text_translate}"
                            print("Used Model Pro Advanced")
                        except Exception as e2:
                            try:
                                response_translation_gen_image = model_translation_gen_prompt_flash.generate_content(check_prompt_response_text)
                                response_translation_gen_image_text_translate = response_translation_gen_image.text.strip()
                                prompt = f"{response_translation_gen_image_text_translate}"
                                print("Used Model Flash")
                            except Exception as e:
                                print(f"Failed to running all Models for translating image generation prompt. | ERROR: {e}")
                                return
                            print(f"Failed running Model Pro Advanced, Running Model Flash | ERROR: {e2}")
                        print(f"Failed running Model Pro, Running Model Pro Advanced | ERROR: {e3}")
                    print("Translated prompt!")
                elif "5" in response_check_text_check:
                    if safegen:
                        print(f"Inappropriate image generation prompt at {interaction.channel.mention} | {prompt}")
                        error_message = "I'm unable to create an image based on your request. Please make sure your prompt aligns with our image generation guidelines."
                        await interaction.followup.send(error_message)
                        add_to_history("Error", error_message)

                        if create_mod_channel:
                            if not mod_channel_name:
                                mod_channel_name = "🔧・mod"
                            mod_channel = discord.utils.get(interaction.guild.channels, name=mod_channel_name)

                            if not mod_channel:
                                mod_channel = await interaction.guild.create_text_channel(mod_channel_name)
                                print(f"Created mod channel {mod_channel_name} for prompts moderation.")
                            
                            if mod_channel:
                                # Create an embed for mod alert
                                embed = discord.Embed(
                                    title="⚠️ Inappropriate Image Generation Prompt Flagged",
                                    description="A user's prompt has been flagged for moderation:",
                                    color=discord.Color.red(),
                                )
                                embed.add_field(name="User", value=f"{interaction.user.mention} (ID: {interaction.user.id})", inline=False)
                                embed.add_field(name="Channel", value=interaction.channel.mention, inline=True)
                                embed.add_field(name="Prompt", value=prompt, inline=False)
                                embed.set_thumbnail(url=interaction.user.avatar.url)
                                embed.set_footer(text="⚙️ Automated Moderation System")

                                # Send the embed with moderation action buttons
                                message = await mod_channel.send(
                                    embed=embed,
                                    view=ModActionView(interaction.user, mod_channel, interaction.channel)
                                )
                        return
                    is_nsfw = True
                elif "4" in response_check_text_check:
                    print("Prompt is safe for image generation.")
                    prompt = check_prompt_response_text
                elif "3" in response_check_text_check:
                    print("Corrupted image generation prompt")
                    add_to_history(
                        "Error",
                        "Oops, something seems off with that prompt. Please try rephrasing it or using different keywords. I'm here to help if you need suggestions!",
                    )
                    await message.channel.send(
                        "Oops, something seems off with that prompt. Please try rephrasing it or using different keywords. I'm here to help if you need suggestions!"
                    )
                    return
                elif "2" in response_check_text_check:
                    if safegen:
                        print(f"Inappropriate image generation prompt at {interaction.channel.mention} | {prompt}")
                        error_message = "I'm unable to create an image based on your request. Please make sure your prompt aligns with our image generation guidelines."
                        await interaction.followup.send(error_message)
                        add_to_history("Error", error_message)
                        return
                    is_nsfw = True
                else:
                    print(f"Checking prompt Error: {response_check_text_check} isn't an available option")
                    add_to_history("System", "Oops! Looks like the image generator took a coffee break ☕. Please try again in a moment!")
                    await interaction.followup.send("Oops! Looks like the image generator took a coffee break ☕. Please try again in a moment!")
                    return
            except Exception as e:
                error_message = str(e)
                if "Unrecognized BlockReason enum value" in error_message or "Invalid operation" in error_message or "cannot access local variable" in error_message:
                    if safegen:
                        print(f"Inappropriate image generation prompt at {interaction.channel.mention} | {prompt}")
                        error_message = "I'm unable to create an image based on your request. Please make sure your prompt aligns with our image generation guidelines."
                        await interaction.followup.send(error_message)
                        add_to_history("Error", error_message)

                        if create_mod_channel:
                            if not mod_channel_name:
                                mod_channel_name = "🔧・mod"
                            mod_channel = discord.utils.get(interaction.guild.channels, name=mod_channel_name)

                            if not mod_channel:
                                mod_channel = await interaction.guild.create_text_channel(mod_channel_name)
                                print(f"Created mod channel {mod_channel_name} for prompts moderation.")
                            
                            if mod_channel:
                                # Create an embed for mod alert
                                embed = discord.Embed(
                                    title="⚠️ Inappropriate Image Generation Prompt Flagged",
                                    description="A user's prompt has been flagged for moderation:",
                                    color=discord.Color.red(),
                                )
                                embed.add_field(name="User", value=f"{interaction.user.mention} (ID: {interaction.user.id})", inline=False)
                                embed.add_field(name="Channel", value=interaction.channel.mention, inline=True)
                                embed.add_field(name="Prompt", value=prompt, inline=False)
                                embed.set_thumbnail(url=interaction.user.avatar.url)
                                embed.set_footer(text="⚙️ Automated Moderation System")

                                # Send the embed with moderation action buttons
                                message = await mod_channel.send(
                                    embed=embed,
                                    view=ModActionView(interaction.user, mod_channel, interaction.channel)
                                )
                        return
                    is_nsfw = True
                else:
                    print(f"Checking prompt Error: {e}")
                    add_to_history("System", "Oops! Looks like the image generator took a coffee break ☕. Please try again in a moment!")
                    await interaction.followup.send("Oops! Looks like the image generator took a coffee break ☕. Please try again in a moment!")
                    return
        
        api_key = HUGGING_FACE_API
        max_retries = 10  # Increased retries for better handling
        backoff_factor = 3  # Increased backoff factor for longer wait times

        member_name = interaction.user.display_name

        # Use the default model if no model is provided
        if model is None:
            model = f"{Image_Model}"

        if model == "stabilityai/stable-diffusion-xl-base-1.0":
            model_name = "Stable Diffusion XL Base 1.0"
        elif model == f"{Image_Model}":
            model_name = f"{Image_Model_Name}"
        elif model == "ehristoforu/dalle-3-xl-v2":
            model_name = "DALL-E 3 XL V2"
        elif model == "black-forest-labs/FLUX.1-schnell":
            model_name = "FLUX.1 Schnell"
        elif model == "dataautogpt3/FLUX-anime2":
            model_name = "FLUX Anime 2"
        elif model == "Yntec/Chip_n_DallE":
            model_name = "Chip & DallE"
        elif model == "black-forest-labs/FLUX.1-dev":
            model_name = "Flux.1 DEV"
        elif model == "stabilityai/stable-diffusion-3-medium-diffusers":
            model_name = "Stable Diffusion 3 Medium Diffusers"
        elif model == "Shakker-Labs/AWPortrait-FL":
            model_name = "AWPortrait FL"
        elif model == "Shakker-Labs/FLUX.1-dev-LoRA-Garbage-Bag-Art":
            model_name = "FLUX.1 DEV LoRA Art"
        elif model == "Shakker-Labs/FLUX.1-dev-LoRA-playful-metropolis":
            model_name = "Flux.1 DEV LoRA Playful Metropolis Art"
        elif model == "Shakker-Labs/FLUX.1-dev-LoRA-Logo-Design":
            model_name = "Flux.1 DEV LoRA Logo Design"
        elif model == "Shakker-Labs/FLUX.1-dev-LoRA-add-details":
            model_name = "Flux.1 DEV LoRA Add Details"
        else:
            model_name = "Unkown Model"

        add_to_history(member_name, f"/img {prompt} | Model: {model_name}")

        url = f'https://api-inference.huggingface.co/models/{model}'
        headers = {
            'Authorization': f'Bearer {api_key}'
        }
        data = {
            'inputs': prompt
        }

        def save_image(response):
            image_path = "system/RAM/gen-image/generated_image.png"
            os.makedirs(os.path.dirname(image_path), exist_ok=True)
            with open(image_path, 'wb') as f:
                f.write(response.content)
            print("Image saved successfully as 'generated_image.png'!")

        async def handle_error(response):
            error_message = response.json().get('error', 'No error message')
            if response.status_code == 503:
                print(f"Service unavailable. Error: {error_message}")
                await interaction.followup.send("Oopsies! Looks like our generator engine are taking a little snooze! Please try again later, and maybe bring some coffee")
            elif response.status_code == 429:
                print(f"Rate limit exceeded. Error: {error_message}")
                await interaction.followup.send("Uh-oh! We've encountered a Rate Limit Error! Please try again later.")
            elif response.status_code == 500:
                print(f"Internal Server Error. | Error: {error_message}")
                await interaction.followup.send("Uh-oh! We've encountered an Internal Server Error! Our processing engine is having a little meltdown. Please try again shortly.")
            else:
                print(f"Failed to save image. Status code: {response.status_code}, Error: {error_message}")

        def fetch_image_with_retries(url, headers, data):
            for attempt in range(max_retries):
                response = requests.post(url, headers=headers, json=data)
                if response.ok:
                    save_image(response)
                    return True
                else:
                    handle_error(response)
                    if response.status_code in [503, 429]:
                        wait_time = backoff_factor ** attempt
                        print(f"Retrying in {wait_time} seconds...")
                        time.sleep(wait_time)
                    else:
                        break
            print("Exceeded maximum retries or encountered a non-retryable error.")
            return False

        success = False
        if model in ["ehristoforu/dalle-3-xl-v2", "black-forest-labs/FLUX.1-schnell", "dataautogpt3/FLUX-anime2", "Shakker-Labs/AWPortrait-FL"]:
            success = fetch_image_with_retries(url, headers, data)
        else:
            response = requests.post(url, headers=headers, json=data)
            if response.ok:
                save_image(response)
                success = True
            else:
                handle_error(response)

        if success:
            image_path = "system/RAM/gen-image/generated_image.png"
            file_extension = image_path.split('.')[-1].lower()
            if file_extension == 'jpg':
                file_extension = 'jpeg'
            if add_watermark_to_generated_image:
                add_watermark("system/RAM/gen-image/generated_image.png", "system/RAM/gen-image/generated_image.png")
            file_path = os.path.join('system/RAM/read-img', f'image.{file_extension}')

            try:
                img = Image.open(image_path).convert('RGB') if file_extension == 'jpeg' else Image.open(image_path)
                buffered = io.BytesIO()
                img.save(buffered, format="PNG")
                img_bytes = buffered.getvalue()
                if not is_nsfw:
                    response = model_V3.generate_content(img)  # Using the original language model
                    analysis_result = response.text.strip()
                    print(f"Image analysis: {analysis_result}")

                    add_to_history_bot("Generated_image", analysis_result)

            except Exception as e:

                print(f"Error: {e}")
                analysis_result = "Error."
                add_to_history("System", f"Error: {str(e)}")
            if is_nsfw:
                try:
                    embed = discord.Embed(title="Generated Image! 🔞 NSFW ⚠️",
                    description=f"⚠️ **This image is marked as NSFW. View at your own discretion.**\n\nPrompt: ||{prompt}||\n\n",
                    color=embed_colors)
                    file = discord.File(image_path, filename="generated_image.png")
                    file.filename = f"SPOILER_{file.filename}"
                    embed.set_image(url="attachment://generated_image.png")
                    embed.set_footer(text=f"Generated by {interaction.user.display_name}\nModel: {model_name}")
                    await interaction.followup.send(embed=embed, file=file)
                except Exception as e:
                    print(f"Error: {e}")
            else:
                embed = discord.Embed(title="Generated Image!",
                    description=f"{prompt}\n",
                    color=embed_colors)
                file = discord.File(image_path, filename="generated_image.png")
                embed.set_image(url="attachment://generated_image.png")
                embed.set_footer(text=f"Generated by {interaction.user.display_name}\nModel: {model_name}")
                await interaction.followup.send(file=file, embed=embed)

            os.remove(image_path)

        else:
            add_to_history("System", "Failed to generate the image after retries.")
            await interaction.followup.send("An error occurred while generating the image. Please try again later or select a different model.")

if DEFAULT_MUSIC_MODEL == "facebook/musicgen-small":
    def_music_model_name = "MusicGen Small"
else:
    def_music_model_name = DEFAULT_MUSIC_MODEL

# Define the list of models as choices
music_model_choices = [
    app_commands.Choice(name="MusicGen Stereo Small", value="facebook/musicgen-stereo-small"),
    app_commands.Choice(name=f"{def_music_model_name} (Default)", value=f"{DEFAULT_MUSIC_MODEL}")
]

# Define the music generation command
@bot.tree.command(name="music", description="Generate music based on your prompt.")
@app_commands.describe(prompt="The prompt for generating the music", model="Choose a model for generating the music (optional)")
@app_commands.choices(model=music_model_choices)
async def generate_music(interaction: discord.Interaction, prompt: str, model: str = "facebook/musicgen-small"):
    if HUGGING_FACE_API == "YOUR_HUGGING_FACE_API_KEY":
        await interaction.followup.send("Sorry, You have entered an Invalid Hugging Face API Key!") 
        return

    await interaction.response.defer()  # Defer the response to allow for processing time
    member_name = interaction.user.display_name

    api_key = HUGGING_FACE_API
    max_retries = 10
    backoff_factor = 3

    if model == "facebook/musicgen-small":
        model_name = "MusicGen Small"  # Default model name
    elif model == "facebook/musicgen-stereo-small":
        model_name = "MusicGen Stereo Small"
    else:
        model_name = model

    add_to_history(member_name, f"/music {prompt} | Model: {model_name}")
    print(f"Using model: {model_name}")
    url = f'https://api-inference.huggingface.co/models/{model}'
    headers = {'Authorization': f'Bearer {api_key}'}
    data = {'inputs': prompt}

    def save_audio(response):
        audio_dir = "system/RAM/gen-music"
        os.makedirs(audio_dir, exist_ok=True)
        audio_path = os.path.join(audio_dir, "generated_music.wav")
        with open(audio_path, 'wb') as f:
            f.write(response.content)
        logging.info(f"Audio generated and saved successfully as '{audio_path}'!")

    def handle_error(response):
        error_message = response.json().get('error', 'No error message')
        if response.status_code == 503:
            logging.error(f"Service unavailable. Error: {error_message}")
        elif response.status_code == 429:
            logging.error(f"Rate limit exceeded. Error: {error_message}")
        else:
            logging.error(f"Failed to generate/save audio. Status code: {response.status_code}, Error: {error_message}")

    def fetch_audio_with_retries(url, headers, data):
        for attempt in range(max_retries):
            response = requests.post(url, headers=headers, json=data)
            if response.ok:
                save_audio(response)
                return True
            else:
                handle_error(response)
                if response.status_code in [503, 429]:
                    wait_time = backoff_factor ** attempt
                    logging.info(f"Retrying in {wait_time} seconds...")
                    time.sleep(wait_time)
                else:
                    break
        logging.error("Exceeded maximum retries or encountered a non-retryable error.")
        return False

    success = fetch_audio_with_retries(url, headers, data)
    if success:
        audio_path = "system/RAM/gen-music/generated_music.wav"
        file = discord.File(audio_path, filename="generated_music.wav")
        await interaction.followup.send(file=file)
        os.remove(audio_path)
    else:
        await interaction.followup.send("An error occurred while generating the music. Please try again later.")

# Supported Models
model_choices = [
    app_commands.Choice(name="Gemini 1.5 Flash | Quick responses and reliable performance.", value="gemini-1.5-flash"),
    app_commands.Choice(name="Gemini 1.5 Flash 002 | Smarter and more reliable than 1.5 Flash for quick and precise responses.", value="gemini-1.5-flash-002"),
    app_commands.Choice(name="Gemini 1.5 Flash (Latest) | Latest Flash advancements, ideal for testing.", value="gemini-1.5-flash-latest"),
    app_commands.Choice(name="Gemini 1.5 Flash 8B | Rapid output for simple prompts and quick interactions.", value="gemini-1.5-flash-8b"),
    app_commands.Choice(name="Gemini 1.5 Pro | Superior depth and understanding for complex tasks.", value="gemini-1.5-pro"),
    app_commands.Choice(name="Gemini 1.5 Pro 002 | Enhanced Gemini 1.5 Pro with superior accuracy and advanced task understanding.", value="gemini-1.5-pro-002"),
    app_commands.Choice(name="Gemini 1.5 Pro (Latest) | The Latest Version of Gemini 1.5 Pro, ideal for testing.", value="gemini-1.5-pro-latest"),
    app_commands.Choice(name="LearnLM 1.5 Pro (Exp) | AI Tutor: Cutting-edge learning to help you study better & faster.", value="learnlm-1.5-pro-experimental"),
    app_commands.Choice(name="Gemini Experimental 1114 | Google's third-most advanced model, handles many complex tasks.", value="gemini-exp-1114"),
    app_commands.Choice(name="Gemini Experimental 1121 | Google's Second advanced Model, Made for complex reasoning and tasks.", value="gemini-exp-1121"),
    app_commands.Choice(name="Gemini Experimental 1206 | Google's Ultimate advanced Model, Outperforming OpenAI 4o and o1 Preview.", value="gemini-exp-1206"),
]

@bot.tree.command(name='model', description="Change your preferred model.")
@app_commands.describe(model="Select a model from the list.")
@app_commands.choices(model=model_choices)
async def change_model(interaction: discord.Interaction, model: str):
    await interaction.response.defer()
    """Change the user's preferred model."""
    user = interaction.user.name
    set_user_model(user, model)
    if model == "gemini-1.5-flash":
        model_name = "Gemini 1.5 Flash"
    elif model == "gemini-1.5-flash-latest":
        model_name = "Gemini 1.5 Flash Latest"
    elif model == "gemini-1.5-flash-8b":
        model_name = "Gemini 1.5 Flash 8B"
    elif model == "gemini-1.5-pro":
        model_name = "Gemini 1.5 Pro"
    elif model == "gemini-1.5-pro-latest":
        model_name = "Gemini 1.5 Pro Latest"
    elif model == "gemini-exp-1114":
        model_name = "Gemini Experimental 1114"
    elif model == "gemini-exp-1121":
        model_name = "Gemini Experimental 1121"
    elif model == "learnlm-1.5-pro-experimental":
        model_name = "LearnLM 1.5 Pro"
    elif model == "gemini-1.5-flash-002":
        model_name = "Gemini 1.5 Flash 002"
    elif model == "gemini-1.5-pro-002":
        model_name = "Gemini 1.5 Pro 002"
    elif model == "gemini-exp-1206":
        model_name = "Gemini Experimental 1206"
    else:
        model_name = model
    await interaction.followup.send(f"Updated model to {model_name}!")

# Supported languages
language_choices = [
    app_commands.Choice(name="English (Default)", value="en"),
    app_commands.Choice(name="Russian", value="ru"),
    app_commands.Choice(name="Spanish", value="es"),
    app_commands.Choice(name="French", value="fr"),
    app_commands.Choice(name="German", value="de"),
    app_commands.Choice(name="Arabic (Egypt)", value="eg"),
    app_commands.Choice(name="Arabic", value="ar"),
]

@bot.tree.command(name="lang", description="Change language for the bot. (Experimental)")
@app_commands.describe(lang="Choose the language")
@app_commands.choices(lang=language_choices)  # Attach the language choices as options
async def change_lang(interaction: discord.Interaction, lang: str):
    global default_lang
    global VOICES
    global model, model_file_flash, model_file_a, model_file, model_vid_a, model_vid, model_V3, model_V2, model_V, model_pro, model_flash, EN_insV2, EN_file_ins, EN_insV, EN_video_ins, EN_ins  # Access the global model variable
    global ins, gen_config, sys_security, genai_model, insV, insV2, file_ins, video_ins, chat_session

    if not lang:  # Check for empty string
        await interaction.response.send_message("Please provide a language to change to.")
        return

    try:
        if lang == "en" or lang == "english":
            ins = EN_ins
            if fix_repeating_prompts:
                ins = f"{ins}\n{fix_mem_ins}"
            video_ins = EN_video_ins
            insV = EN_insV
            file_ins = EN_file_ins
            insV2 = EN_insV2
            VOICES = [
                'en-US-BrianNeural', 'en-US-JennyNeural', 'en-US-GuyNeural', 'en-GB-SoniaNeural', 
                'en-AU-NatashaNeural', 'en-IN-NeerjaNeural', 'en-NZ-MitchellNeural', 'en-CA-ClaraNeural', 
                'en-IE-EmilyNeural', 'en-SG-WayneNeural', 'en-ZA-LeonNeural', 'en-GB-RyanNeural',
                'en-AU-WilliamNeural', 'en-IN-PrabhatNeural', 'en-NZ-MollyNeural', 'en-CA-LiamNeural', 
                'en-IE-OrlaNeural', 'en-SG-LunaNeural', 'en-US-AriaNeural', 'en-GB-MaisieNeural'
            ]
            default_lang = "en"
            await interaction.response.send_message("Successfully changed language to English!")
        
        elif lang == "eg" or lang == "egypt":
            ins = eg_ar_ins
            video_ins = EN_video_ins
            insV = EN_insV
            file_ins = EN_file_ins
            insV2 = EN_insV2
            if fix_repeating_prompts:
                ins = f"{ins}\n{eg_fix_mem_ins}"
            VOICES = ["ar-EG-ShakirNeural", "ar-EG-SalmaNeural"]
            default_lang = "eg"
            await interaction.response.send_message("تم تغيير اللغة إلى العربية المصرية بنجاح!")

        elif lang == "ar" or lang == "arabic":
            ins = ins_ar
            video_ins = EN_video_ins
            insV = EN_insV
            file_ins = EN_file_ins
            insV2 = EN_insV2
            if fix_repeating_prompts:
                ins = f"{ins}\n{ar_fix_mem_ins}"
            VOICES = ["ar-EG-ShakirNeural", "ar-EG-SalmaNeural"]  # Adjust this based on available Arabic voices
            default_lang = "ar"
            await interaction.response.send_message("تم تغيير اللغة إلى العربية بنجاح!")

        elif lang == "ru" or lang == "russian":
            ins = ru_ins
            video_ins = ru_video_ins
            insV = ru_insV
            file_ins = ru_file_ins
            insV2 = ru_insV2
            if fix_repeating_prompts:
                ins = f"{ins}\n{ru_fix_mem_ins}"
            VOICES = ["ru-RU-DmitryNeural", "ru-RU-SvetlanaNeural"]
            default_lang = "ru"
            await interaction.response.send_message("Язык успешно изменен на русский!")

        elif lang == "es" or lang == "spanish":
            ins = es_ins
            video_ins = es_video_ins
            insV = EN_insV
            file_ins = es_file_ins
            insV2 = EN_insV2
            if fix_repeating_prompts:
                ins = f"{ins}\n{es_fix_mem_ins}"
            VOICES = ["es-ES-HelenaNeural", "es-ES-PabloNeural"]
            default_lang = "es"
            await interaction.response.send_message("¡El idioma se ha cambiado correctamente a español!")

        elif lang == "fr" or lang == "french":
            ins = fr_ins
            video_ins = fr_video_ins
            insV = fr_insV
            file_ins = fr_file_ins
            insV2 = fr_insV2
            if fix_repeating_prompts:
                ins = f"{ins}\n{fr_fix_mem_ins}"
            VOICES = ["fr-FR-DeniseNeural", "fr-FR-MathieuNeural"]
            default_lang = "fr"
            await interaction.response.send_message("La langue a été changée avec succès en français!")

        elif lang == "de" or lang == "german":
            ins = de_ins
            video_ins = de_video_ins
            insV = de_insV
            file_ins = de_file_ins
            insV2 = de_insV2
            if fix_repeating_prompts:
                ins = f"{ins}\n{de_fix_mem_ins}"
            VOICES = ["de-DE-KatjaNeural", "de-DE-MichaelNeural"]
            default_lang = "de"
            await interaction.response.send_message("Sprache erfolgreich auf Deutsch geändert!")

        else:
            await interaction.response.send_message(f"Sorry, `{lang}` isn't supported yet.")
            return

        # Reinitialize models based on the selected language
        model = genai.GenerativeModel(
            model_name=genai_model,
            generation_config=gen_config,
            system_instruction=(ins),
            safety_settings=sys_security
        )
        
        model_flash = genai.GenerativeModel( 
            model_name="gemini-1.5-flash",
            generation_config=gen_config,
            system_instruction=(ins),
            safety_settings=sys_security
        )

        model_pro = genai.GenerativeModel( 
            model_name="gemini-1.5-pro-latest",
            generation_config=gen_config,
            system_instruction=(insV),
            safety_settings=sys_security
        )

        model_V = genai.GenerativeModel( 
            model_name=advanced_model,
            generation_config=gen_config,
            system_instruction=(insV),
            safety_settings=sys_security
        )

        model_V2 = genai.GenerativeModel( 
            model_name="gemini-1.5-flash",
            generation_config=gen_config,
            system_instruction=(insV),
            safety_settings=sys_security
        )

        model_V3 = genai.GenerativeModel( 
            model_name="gemini-1.5-flash",
            generation_config=gen_config,
            system_instruction=(insV2),
            safety_settings=sys_security
        )

        model_vid = genai.GenerativeModel(
            model_name="gemini-1.5-pro-latest",
            generation_config=gen_config,
            system_instruction=(video_ins),
        )

        model_vid_a = genai.GenerativeModel(
            model_name=advanced_model,
            generation_config=gen_config,
            system_instruction=(video_ins),
        )

        model_file = genai.GenerativeModel(
            model_name="gemini-1.5-pro-latest",
            generation_config=gen_config,
            system_instruction=(file_ins),
        )

        model_file_a = genai.GenerativeModel(
            model_name=advanced_model,
            generation_config=gen_config,
            system_instruction=(file_ins),
        )

        model_file_flash = genai.GenerativeModel(
            model_name="gemini-1.5-flash",
            generation_config=gen_config,
            system_instruction=(file_ins),
        )

        chat_session = model.start_chat()

    except Exception as e:
        print(f"Error: {e} | {str(e)}")
        await interaction.response.send_message(f"An error occurred: {str(e)}")
        add_to_history("Error occurred", str(e))

# Define the tts command
@bot.tree.command(name="tts", description="Text-to-speech conversion.")
@app_commands.describe(text="The text for converting to voice")
async def tts(interaction: discord.Interaction, text: str):
    await interaction.response.defer()
    global skip_ffmpeg_check
    os.makedirs("system/RAM/vc", exist_ok=True)
    text_content = text

    if text_content:
        
        try:
            skip_ffmpeg_check = True
            await generate_tts(interaction, text_content)
            skip_ffmpeg_check = False

            voice_path = "system/RAM/vc/Generated_voice.mp3"
            voice_path_wav = "system/RAM/vc/Generated_voice.wav"

            # Check if the generated file exists
            file_path = voice_path if os.path.exists(voice_path) else voice_path_wav if os.path.exists(voice_path_wav) else None
            
            if file_path:
                file = discord.File(file_path, filename=os.path.basename(file_path))
                
                await interaction.followup.send(file=file)
                
                # Remove the generated file after use
                os.remove(file_path)
            else:
                await interaction.followup.send("An error occurred generating the TTS file.")
        
        except Exception as e:
            print(f"Error generating TTS: {e}")
            await interaction.followup.send(f"Error generating TTS: {str(e)}")
    else:
        await interaction.followup.send("Please provide some text to convert to voice.", ephemeral=True)

toggle_choices = [
    app_commands.Choice(name="On", value=1),
    app_commands.Choice(name="Off", value=0)
]

@bot.tree.command(name="aitoggle", description="Enable or disable automatic AI responses for this channel.")
@app_commands.describe(toggle="On or off?")
@app_commands.choices(toggle=toggle_choices)
async def aitoggle(interaction: discord.Interaction, toggle: int):
    await interaction.response.defer()
    global ai_toggle_per_channel
    channel_id = interaction.channel_id
    member_name = interaction.user.display_name
    toggle_bool = toggle == 1

    if toggle_bool:
        if channel_id not in ai_toggle_per_channel or not ai_toggle_per_channel[channel_id]:
            ai_toggle_per_channel[channel_id] = True
            await interaction.followup.send(f"Automatic AI responses have been enabled for {interaction.channel.name}.")
            add_to_history(member_name, f"/aitoggle {toggle}")
            add_to_history("System", f"Automatic AI responses have been enabled for {interaction.channel.name}.")
        else:
            await interaction.followup.send(f"Automatic AI responses were already enabled for {interaction.channel.name}.")
            add_to_history(member_name, f"/aitoggle {toggle}")
            add_to_history("System", f"Automatic AI responses were already enabled for {interaction.channel.name}.")
    else:
        if channel_id in ai_toggle_per_channel and ai_toggle_per_channel[channel_id]:
            ai_toggle_per_channel[channel_id] = False
            await interaction.followup.send(f"Automatic AI responses have been disabled for {interaction.channel.name}.")
            add_to_history(member_name, f"/aitoggle {toggle}")
            add_to_history("System", f"Automatic AI responses have been disabled for {interaction.channel.name}.")
        else:
            await interaction.followup.send(f"Automatic AI responses were already disabled for {interaction.channel.name}.")
            add_to_history(member_name, f"/aitoggle {toggle}")
            add_to_history("System", f"Automatic AI responses were already disabled for {interaction.channel.name}.")

bot.run(TOKEN)
